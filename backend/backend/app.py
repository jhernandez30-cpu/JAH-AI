import streamlit as st
import os, tempfile, requests, hashlib, re, json, subprocess, sys, math
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from html import escape
from pathlib import Path
from urllib.parse import urlparse
import sqlite3
import bcrypt

# --- Procesadores de documentos ---
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import yt_dlp
from faster_whisper import WhisperModel
import imageio_ffmpeg
import trafilatura

# --- Chunking y base vectorial ---
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
from chromadb.config import Settings

# --- LLM ---
from langchain_ollama import OllamaLLM

# --- Agency Brain ---
from agency_brain import (
    build_agency_context,
    clear_agency_cache,
    get_agency_status,
    retrieve_agency_agents,
)
from jarvis_brain import (
    build_unified_brain_context,
    get_jarvis_stack_summary,
    get_profile,
    get_programming_profiles,
)
from connected_brain import (
    build_connected_brain_context,
    build_quick_code_docs,
    retrieve_connected_workspace_docs,
)
from programming_skills import get_programming_skill_catalog
from project_workspace import (
    summarize_workspace,
)
from local_model_router import (
    AUTO_MODEL_OPTION,
    choose_local_model,
    get_model_plan,
    sort_models_for_ui,
)
from services.brain_orchestrator import notebooklm_result_to_doc, notebooklm_status_message
from services.brain_connector import BrainConnector
from services.conversation_context import resolve_user_request
from services.notebooklm_service import NotebookLMService
from services.safety_filter import classify_safety, safe_refusal_only_if_really_needed
from services.technical_generators import generate_technical_answer, should_use_technical_generator
from services.technical_intent_router import detect_technical_intent
from services.unified_brain import BrainSourceContext, UnifiedBrain

# ============= CONFIGURACIÓN =============
BACKEND_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = BACKEND_DIR.parent
PERSIST_DIR = os.getenv("TUTOR_IA_PERSIST_DIR", str(PROJECT_ROOT / "vectores" / "brain_db"))
SQLITE_DB = os.getenv("TUTOR_IA_SQLITE_DB", str(PROJECT_ROOT / "database" / "users.db"))
COLLECTION_NAME = os.getenv("TUTOR_IA_COLLECTION", "conocimiento_fast")
LLM_MODEL = os.getenv("TUTOR_IA_LLM_MODEL", "llama3.2:1b")
RECOMMENDED_OLLAMA_MODEL = os.getenv("TUTOR_IA_RECOMMENDED_MODEL", "llama3.2:1b")
OLLAMA_NUM_CTX = int(os.getenv("TUTOR_IA_OLLAMA_NUM_CTX", "4096"))
OLLAMA_NUM_PREDICT = int(os.getenv("TUTOR_IA_OLLAMA_NUM_PREDICT", "3072"))
OLLAMA_TEMPERATURE = float(os.getenv("TUTOR_IA_OLLAMA_TEMPERATURE", "0.2"))
OLLAMA_KEEP_ALIVE = os.getenv("TUTOR_IA_OLLAMA_KEEP_ALIVE", "10m")
WEB_FAST_NUM_PREDICT = int(os.getenv("TUTOR_IA_WEB_FAST_NUM_PREDICT", "2048"))
WHISPER_MODEL = os.getenv("TUTOR_IA_WHISPER_MODEL", "small")
WHISPER_DEVICE = os.getenv("TUTOR_IA_WHISPER_DEVICE", "cpu")
MAX_BATCH_SOURCES = int(os.getenv("TUTOR_IA_MAX_BATCH_SOURCES", "300"))
FILE_EXTRACT_TIMEOUT = int(os.getenv("TUTOR_IA_FILE_EXTRACT_TIMEOUT", "180"))
URL_EXTRACT_TIMEOUT = int(os.getenv("TUTOR_IA_URL_EXTRACT_TIMEOUT", "240"))
CHUNK_SIZE = int(os.getenv("TUTOR_IA_CHUNK_SIZE", "1200"))
CHUNK_OVERLAP = int(os.getenv("TUTOR_IA_CHUNK_OVERLAP", "180"))
EMBED_DIM = int(os.getenv("TUTOR_IA_EMBED_DIM", "384"))
CHROMA_UPSERT_BATCH_SIZE = int(os.getenv("TUTOR_IA_CHROMA_UPSERT_BATCH_SIZE", "256"))
CHROMA_EXISTING_CHECK_BATCH_SIZE = int(os.getenv("TUTOR_IA_CHROMA_EXISTING_CHECK_BATCH_SIZE", "500"))
MAX_FILE_WORKERS = max(1, int(os.getenv("TUTOR_IA_MAX_FILE_WORKERS", "2")))
MAX_URL_WORKERS = max(1, int(os.getenv("TUTOR_IA_MAX_URL_WORKERS", "4")))
RETRIEVE_CANDIDATES = int(os.getenv("TUTOR_IA_RETRIEVE_CANDIDATES", "8"))
RESPONSE_TOP_K = int(os.getenv("TUTOR_IA_RESPONSE_TOP_K", "2"))
MAX_DOC_CONTEXT_CHARS = int(os.getenv("TUTOR_IA_MAX_DOC_CONTEXT_CHARS", "700"))
PROMPT_HISTORY_TURNS = int(os.getenv("TUTOR_IA_PROMPT_HISTORY_TURNS", "3"))
AGENCY_MATCH_LIMIT = int(os.getenv("TUTOR_IA_AGENCY_MATCH_LIMIT", "2"))
AGENCY_CONTEXT_CHARS = int(os.getenv("TUTOR_IA_AGENCY_CONTEXT_CHARS", "3000"))
NOTEBOOKLM_ENABLED_DEFAULT = os.getenv("NOTEBOOKLM_ENABLED", "false").lower() not in {"0", "false", "no", "off", ""}
NOTEBOOKLM_ACTIVE_ID = os.getenv("NOTEBOOKLM_ACTIVE_ID", "").strip()
ALLOWED_GROUP_RE = re.compile(r"^[a-zA-Z0-9_-]{1,32}$")
ALLOWED_USERNAME_RE = re.compile(r"^[a-zA-Z0-9_.-]{3,32}$")
TOKEN_RE = re.compile(r"[a-záéíóúüñ0-9_]+", re.IGNORECASE)
SOURCE_REQUEST_RE = re.compile(
    r"\b(fuente|fuentes|cita|citas|bibliografia|documento|documentos|de donde|origen)\b",
    re.IGNORECASE,
)
LOW_MEMORY_MODEL_PRIORITY = ["llama3.2:1b", "qwen2.5:1.5b", "gemma3:1b", "llama3.2:3b"]
SMART_SEARCH_UNCONFIGURED_MESSAGE = (
    "La Búsqueda inteligente está activada, pero todavía no hay una API de búsqueda web configurada."
)


def smart_web_search(query):
    """
    Función preparada para búsqueda inteligente en la web.
    Aquí se podrá conectar una API real de búsqueda web como Tavily, SerpAPI,
    Brave Search API, Google Custom Search o similar.
    """
    return {
        "enabled": False,
        "message": SMART_SEARCH_UNCONFIGURED_MESSAGE,
        "results": [],
    }


def get_notebooklm_service(enabled=None):
    return NotebookLMService(
        enabled=NOTEBOOKLM_ENABLED_DEFAULT if enabled is None else enabled,
        active_notebook_id=st.session_state.get("notebooklm_active_id", NOTEBOOKLM_ACTIVE_ID),
    )


@st.cache_resource
def get_brain_connector():
    return BrainConnector()


INTERACTION_MODES = {
    "unified": {
        "label": "Cerebro Unificado",
        "description": "Un solo cerebro conectado a fuentes, Obsidian, Agency, OpenJarvis, NotebookLM, Ollama y voz local.",
        "placeholder": "Pregunta, programa, organiza, estudia o planifica desde una sola conversacion...",
        "quick_prompts": [
            "Analiza mi objetivo y dame el siguiente paso mas util con fuentes, metodo y validacion.",
            "Ayudame a entender este tema y conviertelo en un plan practico de aprendizaje.",
            "Revisa esta idea o codigo como un solo cerebro: arquitectura, riesgos, mejoras y pruebas.",
        ],
        "instructions": """
Modo Cerebro Unificado:
- Eres un solo cerebro conectado, no una lista de modos separados.
- Eres un asistente de programacion. Tu prioridad es ayudar a crear, explicar, depurar y estructurar soluciones de software, bases de datos, APIs, interfaces, algoritmos, documentacion tecnica y proyectos academicos.
- Las solicitudes normales de programacion y bases de datos son permitidas y deben resolverse de forma completa.
- Decide internamente si conviene ensenar, organizar, crear, programar, auditar o coordinar especialistas.
- Integra fuentes privadas, Obsidian, Agency, OpenJarvis, NotebookLM, Ollama y voz local como capas del mismo razonamiento.
- Prioriza claridad, accion concreta y validacion.
- Si falta informacion, dilo en una frase y pide el dato minimo necesario.
- No inventes ejecuciones, fuentes ni resultados que no esten en el contexto.
""",
    },
    "study": {
        "label": "Potencia tu estudio",
        "description": "Explica conceptos complejos, usa ejemplos reales y refuerza comprensión.",
        "placeholder": "Pide una explicación simple, ejemplos o una guía de estudio...",
        "quick_prompts": [
            "Explícame los conceptos más difíciles en términos simples y con ejemplos del mundo real.",
            "Crea una guía de estudio con puntos clave, definiciones y preguntas de práctica.",
            "Hazme 5 preguntas para comprobar mi comprensión y luego dame las respuestas.",
        ],
        "instructions": """
Modo Potencia tu estudio:
- Actúa como tutor paciente y claro.
- Explica conceptos complejos en términos simples sin perder precisión.
- Usa ejemplos del mundo real cuando el contexto lo permita.
- Refuerza la comprensión con pasos, analogías, mini-resúmenes o preguntas de práctica.
- Si falta información en el contexto interno, dilo con claridad.
""",
    },
    "organize": {
        "label": "Organiza tu pensamiento",
        "description": "Convierte fuentes en esquemas, argumentos, puntos clave y respaldo.",
        "placeholder": "Pide un esquema, una presentación o una estructura argumental...",
        "quick_prompts": [
            "Crea un esquema de presentación optimizado con puntos clave y evidencia de respaldo.",
            "Organiza este material en introducción, desarrollo, conclusión y posibles preguntas.",
            "Extrae las ideas principales y conviértelas en una estructura clara para exponer.",
        ],
        "instructions": """
Modo Organiza tu pensamiento:
- Actúa como estratega de contenido y pensamiento claro.
- Ordena el material en estructuras útiles: esquema, narrativa, secciones, argumentos y conclusiones.
- Incluye puntos clave y evidencia de respaldo tomada del contexto.
- Ayuda a presentar temas con confianza: anticipa dudas, objeciones y transiciones.
- Si falta evidencia para una afirmación, avisa y sugiere qué fuente haría falta.
""",
    },
    "create": {
        "label": "Elabora nuevas ideas",
        "description": "Detecta tendencias, genera ideas y encuentra oportunidades ocultas.",
        "placeholder": "Pide tendencias, ideas de producto u oportunidades ocultas...",
        "quick_prompts": [
            "Identifica tendencias, patrones y oportunidades ocultas en estas fuentes.",
            "Genera ideas de nuevos productos o mejoras basadas en el material cargado.",
            "Analiza la competencia o el mercado y propón oportunidades accionables.",
        ],
        "instructions": """
Modo Elabora nuevas ideas:
- Actúa como analista creativo y estratégico.
- Identifica patrones, tendencias, tensiones, oportunidades y huecos en el material.
- Genera ideas nuevas conectadas con el contexto, no ocurrencias desconectadas.
- Distingue entre evidencia del contexto, inferencia razonable e hipótesis.
- Propón próximos pasos accionables cuando sea útil.
""",
    },
    "programming": {
        "label": "Cerebro Programador",
        "description": "Orquesta arquitectura, depuracion y revision con disciplina OpenJarvis.",
        "placeholder": "Pide diagnostico de codigo, arquitectura, debugging o revision...",
        "quick_prompts": [
            "Analiza esta idea como arquitecto y convierte el plan en pasos de implementacion verificables.",
            "Haz un diagnostico de bug con ciclo observar, hipotetizar, probar y corregir.",
            "Revisa esta solucion priorizando bugs, seguridad, rendimiento y claridad.",
        ],
        "instructions": """
Modo Cerebro Programador:
- Actua como asistente senior de programacion local-first.
- Usa el perfil Jarvis/OpenJarvis seleccionado como disciplina de razonamiento.
- Para debugging, sigue observar, formular hipotesis, probar y corregir.
- Para arquitectura, respeta el sistema existente, limites claros y cambios pequenos.
- Para review, prioriza correctitud, seguridad, rendimiento y mantenibilidad.
- Si no hay codigo o traceback concreto, no inventes una implementacion completa; entrega checklist y pide el fragmento necesario.
- Entrega acciones concretas, comandos o pruebas solo cuando sean utiles y verificables.
""",
    },
    "agency": {
        "label": "Cerebro Agency",
        "description": "Combina tus fuentes con especialistas de AGENCY-AGENTS-MAIN.",
        "placeholder": "Pide un plan, diagnostico, ruta de agentes o solucion multidisciplinaria...",
        "quick_prompts": [
            "Analiza este objetivo con los especialistas adecuados y dame un plan accionable.",
            "Crea una ruta de agentes para aprender este tema paso a paso con controles de calidad.",
            "Evalua mi idea, detecta riesgos y propon una estrategia de mejora.",
        ],
        "instructions": """
Modo Cerebro Agency:
- Actua como orquestador de especialistas.
- Selecciona el enfoque de los agentes relevantes de Agency segun la pregunta.
- Integra metodologia de expertos con la evidencia recuperada del contexto privado.
- Distingue entre hechos del contexto, criterio experto e inferencias.
- Entrega pasos concretos, criterios de validacion y siguientes acciones cuando sea util.
""",
    },
}

LEGACY_INTERACTION_MODES = {key: value for key, value in INTERACTION_MODES.items() if key != "unified"}
INTERACTION_MODES = {"unified": INTERACTION_MODES["unified"]}

# ============= INICIALIZACIÓN =============
@st.cache_resource(show_spinner="Conectando con Chroma...")
def get_collection():
    client = chromadb.PersistentClient(path=PERSIST_DIR, settings=Settings(anonymized_telemetry=False))
    return client.get_or_create_collection(COLLECTION_NAME)

@st.cache_resource
def get_text_splitter():
    return RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

@st.cache_resource(show_spinner="Conectando con Ollama...")
def get_llm(model_name):
    return OllamaLLM(
        model=model_name,
        temperature=OLLAMA_TEMPERATURE,
        num_ctx=OLLAMA_NUM_CTX,
        num_predict=OLLAMA_NUM_PREDICT,
        keep_alive=OLLAMA_KEEP_ALIVE,
    )


@st.cache_resource(show_spinner="Conectando con Ollama rapido...")
def get_fast_llm(model_name):
    return OllamaLLM(
        model=model_name,
        temperature=OLLAMA_TEMPERATURE,
        num_ctx=OLLAMA_NUM_CTX,
        num_predict=WEB_FAST_NUM_PREDICT,
        keep_alive=OLLAMA_KEEP_ALIVE,
    )


@st.cache_data(ttl=30, show_spinner=False)
def get_installed_ollama_models():
    try:
        completed = subprocess.run(
            ["ollama", "list"],
            capture_output=True,
            text=True,
            timeout=10,
        )
    except Exception:
        return []

    if completed.returncode != 0:
        return []

    models = []
    for line in completed.stdout.splitlines()[1:]:
        parts = line.split()
        if parts:
            models.append(parts[0])
    return models

def choose_llm_model(preferred_model=None, question="", docs=None, brain_context=""):
    models = get_installed_ollama_models()
    return choose_local_model(
        models,
        preferred_model=preferred_model,
        question=question,
        docs=docs,
        brain_context=brain_context,
        fallback_model=LLM_MODEL,
    )

def sort_ollama_models(models):
    return sort_models_for_ui(models)

@st.cache_resource(show_spinner="Cargando Whisper...")
def get_whisper():
    device = WHISPER_DEVICE
    compute_type = "float16" if device == "cuda" else "int8"
    return WhisperModel(WHISPER_MODEL, device=device, compute_type=compute_type)

# ============= BASE DE DATOS DE USUARIOS (SQLITE) =============
def init_users_db():
    with sqlite3.connect(SQLITE_DB) as conn:
        conn.execute('''CREATE TABLE IF NOT EXISTS users
                     (username TEXT PRIMARY KEY, password_hash TEXT, groups TEXT)''')

init_users_db()

def normalize_username(username):
    return (username or "").strip()

def normalize_groups(groups):
    if isinstance(groups, str):
        raw_groups = groups.split(",")
    else:
        raw_groups = groups or []

    clean_groups = []
    for group in raw_groups:
        group = str(group).strip().lower()
        if group and ALLOWED_GROUP_RE.fullmatch(group) and group not in clean_groups:
            clean_groups.append(group)
    return clean_groups or ["public"]

def groups_to_db(groups):
    return ",".join(normalize_groups(groups))

def groups_from_db(groups):
    return normalize_groups(groups)

def get_interaction_mode(mode_key):
    return INTERACTION_MODES.get(mode_key, INTERACTION_MODES["unified"])

def validate_new_user(username, password):
    username = normalize_username(username)
    if not ALLOWED_USERNAME_RE.fullmatch(username):
        return False, "El usuario debe tener 3-32 caracteres: letras, números, punto, guion o guion bajo."
    if len(password or "") < 6:
        return False, "La contraseña debe tener al menos 6 caracteres."
    return True, ""

def count_users():
    with sqlite3.connect(SQLITE_DB) as conn:
        return conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]

def get_all_groups():
    groups = {"public"}
    with sqlite3.connect(SQLITE_DB) as conn:
        for (group_text,) in conn.execute("SELECT groups FROM users"):
            groups.update(normalize_groups(group_text))
    return sorted(groups)

def create_user(username, password, groups="public"):
    username = normalize_username(username)
    ok, message = validate_new_user(username, password)
    if not ok:
        return False, message

    hashed = bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()
    group_text = groups_to_db(groups)
    try:
        with sqlite3.connect(SQLITE_DB) as conn:
            conn.execute("INSERT INTO users (username, password_hash, groups) VALUES (?,?,?)",
                         (username, hashed, group_text))
    except sqlite3.IntegrityError:
        return False, "El usuario ya existe."
    return True, "Usuario creado."

def check_user(username, password):
    username = normalize_username(username)
    if not username or not password:
        return None
    with sqlite3.connect(SQLITE_DB) as conn:
        row = conn.execute("SELECT password_hash, groups FROM users WHERE username=?", (username,)).fetchone()
    if row and bcrypt.checkpw(password.encode(), row[0].encode()):
        return groups_from_db(row[1])
    return None

# ============= FUNCIONES DE EXTRACCIÓN =============
def extract_pdf(file_path):
    chunks = []
    with fitz.open(file_path) as doc:
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                chunks.append({"text": text, "metadata": {"source": file_path, "type": "pdf", "page": i+1}})
    return chunks

def extract_docx(file_path):
    doc = DocxDocument(file_path)
    full = []
    for p in doc.paragraphs:
        if p.text.strip():
            full.append(p.text)
    for table in doc.tables:
        rows = [" | ".join(cell.text for cell in row.cells) for row in table.rows]
        full.append("\n".join(rows))
    text = "\n\n".join(full)
    return [{"text": text, "metadata": {"source": file_path, "type": "docx"}}]

def extract_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    chunks = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(" | ".join(str(c) if c is not None else "" for c in row))
        if rows:
            chunks.append({"text": f"Hoja: {name}\n" + "\n".join(rows),
                           "metadata": {"source": file_path, "type": "excel", "sheet": name}})
    return chunks

def extract_pptx(file_path):
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Nota]: {notes}")
        if texts:
            chunks.append({"text": f"Diapositiva {i+1}\n" + "\n".join(texts),
                           "metadata": {"source": file_path, "type": "pptx", "slide": i+1}})
    return chunks

def extract_video(url):
    with tempfile.TemporaryDirectory() as tmp_dir:
        audio_template = str(Path(tmp_dir) / "audio.%(ext)s")
        ydl_opts = {
            'format': 'bestaudio/best',
            'outtmpl': audio_template,
            'postprocessors': [{'key': 'FFmpegExtractAudio', 'preferredcodec': 'mp3'}],
            'ffmpeg_location': imageio_ffmpeg.get_ffmpeg_exe(),
            'quiet': True,
            'noplaylist': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            title = info.get('title', 'sin título')

        audio_files = list(Path(tmp_dir).glob("*.mp3"))
        if not audio_files:
            raise FileNotFoundError("No se pudo generar el audio para transcribir.")

        model = get_whisper()
        segments, _ = model.transcribe(str(audio_files[0]), language="es")
        chunks = []
        for seg in segments:
            chunks.append({
                "text": f"[{seg.start:.1f}s - {seg.end:.1f}s] {seg.text}",
                "metadata": {"source": url, "type": "video", "title": title, "start": seg.start, "end": seg.end}
            })
    return chunks

def extract_article(url):
    response = requests.get(url, timeout=20, headers={"User-Agent": "TutorIA/1.0"})
    response.raise_for_status()
    html = response.text
    text = trafilatura.extract(html, include_comments=False, include_tables=False)
    if text is None:
        raise ValueError("No se pudo extraer texto.")
    meta = trafilatura.extract_metadata(html)
    title = meta.title if meta and meta.title else url
    return [{"text": text, "metadata": {"source": url, "type": "article", "title": title}}]

def process_manual_text(text, title="Nota manual", tags=None):
    return [{"text": text, "metadata": {"source": "manual", "type": "text", "title": title, "tags": tags or []}}]

def is_valid_url(url):
    parsed = urlparse((url or "").strip())
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)

def parse_url_list(text):
    urls = []
    seen = set()
    for raw in re.split(r"[\r\n]+", text or ""):
        url = raw.strip()
        if not url or url in seen:
            continue
        seen.add(url)
        urls.append(url)
    return urls

def is_video_url(url):
    host = urlparse(url).netloc.lower()
    video_hosts = ("youtube.com", "youtu.be", "tiktok.com", "vimeo.com")
    return any(video_host in host for video_host in video_hosts)

def clean_metadata(metadata):
    clean = {}
    for key, value in metadata.items():
        if value is None:
            continue
        if isinstance(value, (str, int, float, bool)):
            clean[key] = value
        elif isinstance(value, (list, tuple, set)):
            clean[key] = ", ".join(str(item) for item in value if str(item).strip())
        else:
            clean[key] = str(value)
    return clean

def make_doc_id(metadata, text):
    source = metadata.get("source", "unknown")
    location_parts = [
        metadata.get("type", ""),
        metadata.get("page", ""),
        metadata.get("sheet", ""),
        metadata.get("slide", ""),
        metadata.get("start", ""),
        metadata.get("end", ""),
    ]
    location = "|".join(str(part) for part in location_parts)
    digest = hashlib.sha256(f"{source}|{location}|{text}".encode("utf-8")).hexdigest()
    return digest[:32]

def new_memory():
    return []

def add_memory_turn(memory, question, answer, max_turns=12):
    memory.append({"role": "human", "content": question})
    memory.append({"role": "ai", "content": answer})
    max_messages = max_turns * 2
    if len(memory) > max_messages:
        del memory[:-max_messages]

def embed_text(text):
    """Embedding local rápido basado en hashing para evitar modelos pesados en la ingesta."""
    vector = [0.0] * EMBED_DIM
    tokens = TOKEN_RE.findall((text or "").lower())
    if not tokens:
        return vector

    previous = ""
    for token in tokens:
        features = [token]
        if previous:
            features.append(f"{previous}_{token}")
        previous = token

        for feature in features:
            digest = hashlib.blake2b(feature.encode("utf-8"), digest_size=8).digest()
            value = int.from_bytes(digest, "little", signed=False)
            index = value % EMBED_DIM
            vector[index] += 1.0 if value & 1 else -1.0

    norm = math.sqrt(sum(value * value for value in vector))
    if not norm:
        return vector
    return [value / norm for value in vector]

def embed_texts(texts):
    return [embed_text(text) for text in texts]

# ============= CHUNKING E INDEXACIÓN (con grupo de acceso) =============
def prepare_vector_records(raw_chunks, access_group="public"):
    """Trocea los textos y prepara documentos, metadatos e IDs estables para Chroma."""
    text_splitter = get_text_splitter()
    final = []
    for rc in raw_chunks:
        text = (rc.get("text") or "").strip()
        if not text:
            continue
        base_metadata = rc.get("metadata", {}).copy()
        base_metadata["access_group"] = normalize_groups([access_group])[0]
        if len(text) > 1000:
            sub_texts = text_splitter.split_text(text)
            for st in sub_texts:
                final.append({"text": st, "metadata": base_metadata.copy()})
        else:
            final.append({"text": text, "metadata": base_metadata})

    docs, metas, ids = [], [], []
    seen_ids = set()
    for f in final:
        meta = clean_metadata(f["metadata"])
        doc_id = make_doc_id(meta, f["text"])
        if doc_id in seen_ids:
            continue
        seen_ids.add(doc_id)
        docs.append(f["text"])
        metas.append(meta)
        ids.append(doc_id)
    return docs, metas, ids

def filter_existing_records(collection, docs, metas, ids):
    if not ids:
        return docs, metas, ids, 0

    existing_ids = set()
    for start in range(0, len(ids), CHROMA_EXISTING_CHECK_BATCH_SIZE):
        batch_ids = ids[start:start + CHROMA_EXISTING_CHECK_BATCH_SIZE]
        try:
            existing = collection.get(ids=batch_ids)
            existing_ids.update(existing.get("ids", []))
        except Exception:
            continue

    if not existing_ids:
        return docs, metas, ids, 0

    new_docs, new_metas, new_ids = [], [], []
    for doc, meta, doc_id in zip(docs, metas, ids):
        if doc_id in existing_ids:
            continue
        new_docs.append(doc)
        new_metas.append(meta)
        new_ids.append(doc_id)
    return new_docs, new_metas, new_ids, len(existing_ids)

def upsert_records(docs, metas, ids, progress_callback=None):
    if not docs:
        return 0

    collection = get_collection()
    indexed_count = 0
    total = len(docs)
    for start in range(0, total, CHROMA_UPSERT_BATCH_SIZE):
        end = min(start + CHROMA_UPSERT_BATCH_SIZE, total)
        batch_docs = docs[start:end]
        collection.upsert(
            documents=batch_docs,
            metadatas=metas[start:end],
            ids=ids[start:end],
            embeddings=embed_texts(batch_docs),
        )
        indexed_count = end
        if progress_callback:
            progress_callback(indexed_count, total)
    return indexed_count

def add_to_db(raw_chunks, access_group="public", skip_existing=True, progress_callback=None):
    """Trocea, genera embeddings y guarda en Chroma añadiendo access_group al metadata."""
    docs, metas, ids = prepare_vector_records(raw_chunks, access_group)
    prepared_count = len(ids)
    skipped_count = 0
    if docs:
        collection = get_collection()
        if skip_existing:
            docs, metas, ids, skipped_count = filter_existing_records(collection, docs, metas, ids)
        upsert_records(docs, metas, ids, progress_callback)
    return {"indexed": len(docs), "skipped": skipped_count, "prepared": prepared_count}

def extract_uploaded_file_bytes(source_name, file_bytes):
    safe_name = Path(source_name).name
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir) / safe_name
        out_path = Path(tmp_dir) / "extracted.json"
        tmp_path.write_bytes(file_bytes)

        worker_path = Path(__file__).with_name("extract_worker.py")
        try:
            completed = subprocess.run(
                [sys.executable, str(worker_path), str(tmp_path), safe_name, str(out_path)],
                cwd=str(Path(__file__).parent),
                capture_output=True,
                text=True,
                timeout=FILE_EXTRACT_TIMEOUT,
            )
        except subprocess.TimeoutExpired as e:
            raise TimeoutError(f"El archivo tardo mas de {FILE_EXTRACT_TIMEOUT} segundos en extraerse.") from e

        if completed.returncode != 0:
            error_text = (completed.stderr or completed.stdout or "Error desconocido").strip()
            raise RuntimeError(error_text)

        raw = json.loads(out_path.read_text(encoding="utf-8"))
    return safe_name, raw

def extract_uploaded_file(uploaded_file):
    safe_name = Path(uploaded_file.name).name
    return extract_uploaded_file_bytes(safe_name, bytes(uploaded_file.getbuffer()))
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir) / safe_name
        out_path = Path(tmp_dir) / "extracted.json"
        with open(tmp_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        worker_path = Path(__file__).with_name("extract_worker.py")
        try:
            completed = subprocess.run(
                [sys.executable, str(worker_path), str(tmp_path), safe_name, str(out_path)],
                cwd=str(Path(__file__).parent),
                capture_output=True,
                text=True,
                timeout=FILE_EXTRACT_TIMEOUT,
            )
        except subprocess.TimeoutExpired as e:
            raise TimeoutError(f"El archivo tardó más de {FILE_EXTRACT_TIMEOUT} segundos en extraerse.") from e

        if completed.returncode != 0:
            error_text = (completed.stderr or completed.stdout or "Error desconocido").strip()
            raise RuntimeError(error_text)

        raw = json.loads(out_path.read_text(encoding="utf-8"))
    return safe_name, raw

def ingest_uploaded_file(uploaded_file, access_group):
    source_name = Path(uploaded_file.name).name
    try:
        source_name, raw = extract_uploaded_file(uploaded_file)
        if not raw:
            return {"source": source_name, "ok": False, "fragments": 0, "skipped": 0, "error": "No se encontró texto para indexar."}
        stats = add_to_db(raw, access_group)
        return {"source": source_name, "ok": True, "fragments": stats["indexed"], "skipped": stats["skipped"], "error": ""}
    except Exception as e:
        return {"source": source_name, "ok": False, "fragments": 0, "skipped": 0, "error": str(e)}

def ingest_uploaded_files_batch(uploaded_files, access_group, extraction_callback=None, index_callback=None, done_callback=None):
    results = []

    for idx, uploaded_file in enumerate(uploaded_files, start=1):
        source_name = Path(uploaded_file.name).name
        started_at = time.perf_counter()
        if extraction_callback:
            extraction_callback(idx, len(uploaded_files), source_name)
        try:
            source_name, raw = extract_uploaded_file(uploaded_file)
            seconds = time.perf_counter() - started_at
            if not raw:
                results.append({
                    "source": source_name,
                    "ok": False,
                    "fragments": 0,
                    "skipped": 0,
                    "seconds": seconds,
                    "error": "No se encontró texto para indexar.",
                })
                if done_callback:
                    done_callback(idx, len(uploaded_files), source_name)
                continue
            if index_callback:
                index_callback(idx, len(uploaded_files), source_name, 0, 1)
            stats = add_to_db(
                raw,
                access_group,
                progress_callback=lambda done, total, idx=idx, source_name=source_name: index_callback(
                    idx, len(uploaded_files), source_name, done, total
                ) if index_callback else None,
            )
            results.append({
                "source": source_name,
                "ok": True,
                "fragments": stats["indexed"],
                "skipped": stats["skipped"],
                "seconds": time.perf_counter() - started_at,
                "error": "",
            })
        except Exception as e:
            results.append({
                "source": source_name,
                "ok": False,
                "fragments": 0,
                "skipped": 0,
                "seconds": time.perf_counter() - started_at,
                "error": str(e),
            })
        if done_callback:
            done_callback(idx, len(uploaded_files), source_name)

    return results

def ingest_uploaded_files_batch(uploaded_files, access_group, extraction_callback=None, index_callback=None, done_callback=None):
    payloads = [
        (idx, Path(uploaded_file.name).name, bytes(uploaded_file.getbuffer()))
        for idx, uploaded_file in enumerate(uploaded_files, start=1)
    ]
    results_by_index = {}
    total = len(payloads)
    workers = min(MAX_FILE_WORKERS, total) if total else 1

    def index_raw(idx, source_name, raw, started_at):
        if not raw:
            return {
                "source": source_name,
                "ok": False,
                "fragments": 0,
                "skipped": 0,
                "seconds": time.perf_counter() - started_at,
                "error": "No se encontro texto para indexar.",
            }
        if index_callback:
            index_callback(idx, total, source_name, 0, 1)
        stats = add_to_db(
            raw,
            access_group,
            progress_callback=lambda done, total_fragments, idx=idx, source_name=source_name: index_callback(
                idx, total, source_name, done, total_fragments
            ) if index_callback else None,
        )
        return {
            "source": source_name,
            "ok": True,
            "fragments": stats["indexed"],
            "skipped": stats["skipped"],
            "seconds": time.perf_counter() - started_at,
            "error": "",
        }

    if workers <= 1:
        for idx, source_name, file_bytes in payloads:
            started_at = time.perf_counter()
            if extraction_callback:
                extraction_callback(idx, total, source_name)
            try:
                source_name, raw = extract_uploaded_file_bytes(source_name, file_bytes)
                results_by_index[idx] = index_raw(idx, source_name, raw, started_at)
            except Exception as e:
                results_by_index[idx] = {
                    "source": source_name,
                    "ok": False,
                    "fragments": 0,
                    "skipped": 0,
                    "seconds": time.perf_counter() - started_at,
                    "error": str(e),
                }
            if done_callback:
                done_callback(idx, total, source_name)
        return [results_by_index[idx] for idx, _, _ in payloads]

    with ThreadPoolExecutor(max_workers=workers) as executor:
        future_map = {}
        for idx, source_name, file_bytes in payloads:
            if extraction_callback:
                extraction_callback(idx, total, source_name)
            started_at = time.perf_counter()
            future = executor.submit(extract_uploaded_file_bytes, source_name, file_bytes)
            future_map[future] = (idx, source_name, started_at)

        completed_count = 0
        for future in as_completed(future_map):
            idx, source_name, started_at = future_map[future]
            completed_count += 1
            try:
                source_name, raw = future.result()
                results_by_index[idx] = index_raw(idx, source_name, raw, started_at)
            except Exception as e:
                results_by_index[idx] = {
                    "source": source_name,
                    "ok": False,
                    "fragments": 0,
                    "skipped": 0,
                    "seconds": time.perf_counter() - started_at,
                    "error": str(e),
                }
            if done_callback:
                done_callback(completed_count, total, source_name)

    return [results_by_index[idx] for idx, _, _ in payloads]

def ingest_url_sources_batch(urls, access_group, audio_fallback=False, progress_callback=None):
    results_by_index = {}
    total = len(urls)
    workers = 1 if audio_fallback else min(MAX_URL_WORKERS, total or 1)

    def process_url(idx, url):
        started_at = time.perf_counter()
        try:
            if not is_valid_url(url):
                return idx, {
                    "source": url,
                    "ok": False,
                    "fragments": 0,
                    "skipped": 0,
                    "seconds": 0,
                    "error": "URL invalida.",
                }
            raw = extract_url_source(url, audio_fallback=audio_fallback)
            if not raw:
                return idx, {
                    "source": url,
                    "ok": False,
                    "fragments": 0,
                    "skipped": 0,
                    "seconds": time.perf_counter() - started_at,
                    "error": "No se encontro texto para indexar.",
                }
            return idx, {"source": url, "raw": raw, "started_at": started_at}
        except Exception as e:
            return idx, {
                "source": url,
                "ok": False,
                "fragments": 0,
                "skipped": 0,
                "seconds": time.perf_counter() - started_at,
                "error": str(e),
            }

    if workers <= 1:
        completed_payloads = []
        for idx, url in enumerate(urls, start=1):
            if progress_callback:
                progress_callback(idx - 1, total, url, "extrayendo")
            completed_payloads.append(process_url(idx, url))
    else:
        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_map = {}
            for idx, url in enumerate(urls, start=1):
                if progress_callback:
                    progress_callback(idx - 1, total, url, "extrayendo")
                future_map[executor.submit(process_url, idx, url)] = url
            completed_payloads = [future.result() for future in as_completed(future_map)]

    completed_count = 0
    for idx, payload in completed_payloads:
        completed_count += 1
        url = payload["source"]
        if "raw" not in payload:
            results_by_index[idx] = payload
        else:
            if progress_callback:
                progress_callback(completed_count - 1, total, url, "indexando")
            stats = add_to_db(payload["raw"], access_group)
            results_by_index[idx] = {
                "source": url,
                "ok": True,
                "fragments": stats["indexed"],
                "skipped": stats["skipped"],
                "seconds": time.perf_counter() - payload["started_at"],
                "error": "",
            }
        if progress_callback:
            progress_callback(completed_count, total, url, "listo")

    return [results_by_index[idx] for idx in range(1, total + 1)]

def _legacy_ingest_url_source(url, access_group):
    url = url.strip()
    try:
        if not is_valid_url(url):
            return {"source": url, "ok": False, "fragments": 0, "skipped": 0, "error": "URL inválida."}
        raw = extract_video(url) if is_video_url(url) else extract_article(url)
        if not raw:
            return {"source": url, "ok": False, "fragments": 0, "skipped": 0, "error": "No se encontró texto para indexar."}
        stats = add_to_db(raw, access_group)
        return {"source": url, "ok": True, "fragments": stats["indexed"], "skipped": stats["skipped"], "error": ""}
    except Exception as e:
        return {"source": url, "ok": False, "fragments": 0, "skipped": 0, "error": str(e)}

def extract_url_source(url, audio_fallback=False):
    with tempfile.TemporaryDirectory() as tmp_dir:
        out_path = Path(tmp_dir) / "url_extracted.json"
        worker_path = Path(__file__).with_name("url_worker.py")
        worker_env = os.environ.copy()
        worker_env["TUTOR_IA_VIDEO_AUDIO_FALLBACK"] = "1" if audio_fallback else "0"
        try:
            completed = subprocess.run(
                [sys.executable, str(worker_path), url, str(out_path)],
                cwd=str(Path(__file__).parent),
                capture_output=True,
                text=True,
                env=worker_env,
                timeout=URL_EXTRACT_TIMEOUT,
            )
        except subprocess.TimeoutExpired as e:
            raise TimeoutError(f"La URL tardó más de {URL_EXTRACT_TIMEOUT} segundos en procesarse.") from e

        if completed.returncode != 0:
            error_text = (completed.stderr or completed.stdout or "Error desconocido").strip()
            raise RuntimeError(error_text)

        return json.loads(out_path.read_text(encoding="utf-8"))

def ingest_url_source(url, access_group, audio_fallback=False):
    url = url.strip()
    started_at = time.perf_counter()
    try:
        if not is_valid_url(url):
            return {"source": url, "ok": False, "fragments": 0, "skipped": 0, "seconds": 0, "error": "URL inválida."}
        raw = extract_url_source(url, audio_fallback=audio_fallback)
        if not raw:
            return {
                "source": url,
                "ok": False,
                "fragments": 0,
                "skipped": 0,
                "seconds": time.perf_counter() - started_at,
                "error": "No se encontró texto para indexar.",
            }
        stats = add_to_db(raw, access_group)
        return {
            "source": url,
            "ok": True,
            "fragments": stats["indexed"],
            "skipped": stats["skipped"],
            "seconds": time.perf_counter() - started_at,
            "error": "",
        }
    except Exception as e:
        return {
            "source": url,
            "ok": False,
            "fragments": 0,
            "skipped": 0,
            "seconds": time.perf_counter() - started_at,
            "error": str(e),
        }

def render_ingest_results(results):
    if not results:
        return

    ok_results = [r for r in results if r["ok"]]
    failed_results = [r for r in results if not r["ok"]]
    total_fragments = sum(r["fragments"] for r in ok_results)
    total_skipped = sum(r.get("skipped", 0) for r in ok_results)

    if ok_results:
        st.success(f"{len(ok_results)} fuentes procesadas, {total_fragments} fragmentos nuevos.")
        if total_skipped:
            st.info(f"{total_skipped} fragmentos ya existían y se saltaron.")
    else:
        st.error("No se pudo indexar ninguna fuente del lote.")
    if failed_results:
        st.warning(f"{len(failed_results)} fuentes no se pudieron procesar.")
        with st.expander("Ver errores del lote"):
            for result in failed_results:
                st.write(f"- {result['source']}: {result['error']}")

# ============= BÚSQUEDA CON RERANKING Y FILTRO DE ACCESO =============
def retrieve(question, user_groups=None, k=None, top_k=None, selected_sources=None):
    """
    1. Obtiene k candidatos de Chroma con filtro opcional por access_group.
    2. Devuelve los top_k documentos más cercanos.
    """
    collection = get_collection()
    total_docs = collection.count()
    if total_docs == 0:
        return []

    k = k or RETRIEVE_CANDIDATES
    top_k = top_k or RESPONSE_TOP_K
    user_groups = normalize_groups(user_groups or ["public"])
    if selected_sources is not None:
        selected_sources = set(selected_sources)
        if not selected_sources:
            return []
    n_results = min(max(k, top_k * 8), total_docs)
    where_filter = None
    if "admin" not in user_groups:
        if len(user_groups) == 1:
            where_filter = {"access_group": user_groups[0]}
        else:
            where_filter = {"$or": [{"access_group": group} for group in user_groups]}

    try:
        res = collection.query(query_embeddings=[embed_text(question)], n_results=n_results, where=where_filter)
    except Exception:
        res = collection.query(query_embeddings=[embed_text(question)], n_results=n_results)

    docs = []
    if res['documents']:
        for i, doc_text in enumerate(res['documents'][0]):
            meta = res['metadatas'][0][i]
            # Filtro manual por grupo de acceso
            doc_group = meta.get("access_group", "public")
            source = meta.get("source", "")
            source_allowed = selected_sources is None or source in selected_sources
            if source_allowed and (doc_group in user_groups or "admin" in user_groups):
                docs.append({"text": doc_text, "metadata": meta})

    return docs[:top_k]

# ============= GENERACIÓN CON MEMORIA =============
def trim_prompt_text(text, max_chars):
    text = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 3].rstrip() + "..."


def clean_answer_text(text):
    text = str(text or "").replace("**", "")
    text = re.sub(r"(?im)^\s*fuentes?\s*:\s*(?:\n\s*[-*].*)+", "", text)
    text = re.sub(r"(?im)^\s*sources?\s*:\s*(?:\n\s*[-*].*)+", "", text)
    text = re.sub(r"(?im)^\s*fuentes?\s*:\s*.*$", "", text)
    text = re.sub(r"(?im)^\s*sources?\s*:\s*.*$", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def source_requested(question):
    return bool(SOURCE_REQUEST_RE.search(str(question or "")))


def brain_context_metadata(context):
    return getattr(context, "metadata", {}) if not isinstance(context, dict) else context.get("metadata", {})


def brain_context_content(context):
    return getattr(context, "content", "") if not isinstance(context, dict) else context.get("content", "")


def brain_context_by_source(contexts, source):
    return [
        context
        for context in contexts or []
        if (getattr(context, "source", None) if not isinstance(context, dict) else context.get("source")) == source
    ]


def brain_contexts_to_docs(contexts):
    docs = []
    for context in contexts or []:
        docs.extend(brain_context_metadata(context).get("docs", []) or [])
    return docs


def make_context_doc(source, title, text, doc_type="tool"):
    return {
        "text": trim_prompt_text(text, MAX_DOC_CONTEXT_CHARS),
        "metadata": {
            "source": source,
            "title": title,
            "type": doc_type,
            "access_group": "public",
        },
    }


def build_tools_context(question, route):
    if route.get("requires_calculation"):
        match = re.search(r"(-?\d+(?:\.\d+)?)\s*([\+\-\*/])\s*(-?\d+(?:\.\d+)?)", question)
        if match:
            left = float(match.group(1))
            right = float(match.group(3))
            op = match.group(2)
            value = {
                "+": left + right,
                "-": left - right,
                "*": left * right,
                "/": left / right if right != 0 else "division entre cero",
            }[op]
            return f"calculator: {match.group(0)} = {value}"
    if route.get("requires_shell"):
        return "shell_exec disponible: no se ejecutan comandos peligrosos ni destructivos sin confirmacion."
    if route.get("requires_file_read"):
        return "file_read disponible para rutas validas del proyecto."
    if route.get("requires_file_write"):
        return "file_write disponible con validacion de ruta y cambios minimos."
    return ""


def technical_source_name(intent):
    intent_name = str((intent or {}).get("resolved_intent") or (intent or {}).get("intent") or "")
    if intent_name in {"database_design", "sql", "sql_generation", "er_model"}:
        return "database_generator"
    if intent_name in {"code_debugging", "code_review", "code_explanation"}:
        return "code_interpreter_router"
    return "technical_generators"


def try_direct_technical_answer(question, memory=None):
    """
    Fast lane for deterministic technical tasks.
    Avoids sending schema/code generation to Ollama or Bridge when a local
    template can answer immediately.
    """
    technical_intent = detect_technical_intent(question, memory)
    resolved_question = technical_intent.get("resolved_request") or resolve_user_request(question, memory)
    safety = classify_safety(resolved_question, technical_intent, memory)

    if technical_intent.get("needs_clarification"):
        answer = "Te refieres a la base de datos, una pagina, un CRUD u otra estructura?"
        source_name = "conversation_resolver"
    elif not safety.get("allowed", True):
        answer = safe_refusal_only_if_really_needed(str(safety.get("reason") or ""))
        source_name = "safety_filter"
    elif should_use_technical_generator(technical_intent):
        answer = generate_technical_answer(resolved_question, technical_intent)
        source_name = technical_source_name(technical_intent)
    else:
        return None

    answer = clean_answer_text(answer)
    if memory is not None:
        add_memory_turn(memory, question, answer)

    return {
        "answer": answer,
        "sources": [
            {
                "text": trim_prompt_text(answer, 260),
                "metadata": {
                    "source": source_name,
                    "title": source_name,
                    "type": "technical_template",
                    "access_group": "public",
                },
            }
        ],
        "agency_agents": [],
        "brain_parts": ["technical_intent_router", "safety_filter", source_name],
        "notebooklm_used_count": 0,
        "notebooklm_message": "",
        "route": {"technical_intent": technical_intent},
        "source_contexts": [],
        "mode": get_interaction_mode("unified")["label"],
        "sources_used": [source_name],
        "success": True,
        "cached": False,
        "latency_ms": 0,
        "intent": technical_intent,
        "safety": safety,
        "resolved_message": resolved_question,
    }


def generate_answer(question, docs, memory=None, interaction_mode="unified", model_name=None, agency_context="", brain_context="", show_sources=False, fast_response=False):
    mode = get_interaction_mode(interaction_mode)
    model_name = choose_llm_model(
        model_name,
        question=question,
        docs=docs,
        brain_context=brain_context,
    )
    if not model_name:
        response = (
            "No hay modelos de Ollama instalados todavía. "
            f"Descarga uno con: `ollama pull {RECOMMENDED_OLLAMA_MODEL}`. "
            "Después vuelve a intentarlo."
        )
        response = clean_answer_text(response)
        if memory is not None:
            add_memory_turn(memory, question, response)
        return response

    if not docs and not agency_context and not brain_context:
        response = clean_answer_text("No encontre informacion relevante en la base de conocimiento para responder esa pregunta.")
        if memory is not None:
            add_memory_turn(memory, question, response)
        return response

    context = ""
    for d in docs or []:
        meta = d["metadata"]
        tipo = {"video": "🎬 Vídeo", "article": "📄 Artículo", "pdf": "📕 PDF",
                "docx": "📝 Word", "excel": "📊 Excel", "pptx": "🎞️ PPT", "text": "✍️ Nota"}.get(meta["type"], "")
        title = meta.get("title", meta.get("source", "fuente"))
        context += f"[{tipo} {title}]\n{trim_prompt_text(d['text'], MAX_DOC_CONTEXT_CHARS)}\n\n"
    if not context:
        context = "No se recuperaron fuentes privadas relevantes para esta pregunta.\n"

    agency_section = ""
    if agency_context:
        agency_section = f"""
Base Agency:
{agency_context}

Reglas para usar Agency:
- Usa Agency como metodologia interna y como apoyo experto.
- Prioriza el contexto privado cuando exista.
- Si solo Agency respalda una recomendacion, dilo como criterio experto o inferencia.
- No inventes que un especialista ejecuto acciones reales; solo usa su enfoque.
"""

    brain_section = ""
    if brain_context:
        brain_section = f"""
Capa Jarvis/OpenJarvis:
{brain_context}
"""

    # Construir prompt con historial
    history_text = ""
    if memory:
        recent_memory = memory[-(PROMPT_HISTORY_TURNS * 2):]
        for msg in recent_memory:
            if msg.get("role") == "human":
                history_text += f"Estudiante: {trim_prompt_text(msg.get('content', ''), 500)}\n"
            elif msg.get("role") == "ai":
                history_text += f"Tutor: {trim_prompt_text(msg.get('content', ''), 900)}\n"
        if history_text:
            history_text = "Historial de la conversación:\n" + history_text + "\n"

    show_sources = show_sources or source_requested(question)
    source_rule = (
        "- Si el usuario pidio fuentes, menciona solo los titulos estrictamente necesarios al final.\n"
        if show_sources
        else "- No cites fuentes, no agregues secciones de fuentes y no digas 'segun el documento'; usa el contexto en silencio para construir una respuesta inteligente.\n"
    )

    prompt = f"""
Eres el cerebro de una aplicación y página web tipo NotebookLM, conectado a una base de conocimiento privada.
Tu modo actual es: {mode["label"]}.

Reglas generales:
- Eres un asistente de programacion. Tu prioridad es ayudar a crear, explicar, depurar y estructurar soluciones de software, bases de datos, APIs, interfaces, algoritmos, documentacion tecnica y proyectos academicos.
- Las solicitudes normales de programacion y bases de datos son permitidas y deben resolverse de forma completa.
- Cuando el usuario pida crear, hacer, generar, construir o dame, entrega el resultado completo: codigo, SQL, tablas, relaciones, pasos o estructura segun corresponda.
- Si una fuente del Cerebro Unificado no devuelve informacion util, no inventes un rechazo; usa conocimiento general tecnico y responde con una solucion profesional.
- Lee y sintetiza el contexto como material interno; no copies fragmentos largos.
- Responde como tutor técnico inteligente: directo, claro, práctico y con criterio.
- Usa principalmente la información proporcionada en el contexto, pero intégrala con razonamiento técnico.
- Si falta información para una respuesta segura, dilo en una frase y da el mejor siguiente paso.
- Puedes hacer inferencias útiles, pero márcalas como inferencias cuando no estén explícitas en el contexto.
- Responde en español claro.
- No uses negritas Markdown, no escribas ** y evita adornos innecesarios.
- Prioriza respuestas completas en tareas tecnicas; solo se breve cuando el usuario pida una explicacion simple.
{source_rule}

{mode["instructions"]}

{agency_section}
{brain_section}

{history_text}
Contexto:
{context}

Pregunta del estudiante: {question}
Respuesta del tutor:
"""
    try:
        llm = get_fast_llm(model_name) if fast_response else get_llm(model_name)
        response = llm.invoke(prompt)
    except Exception as e:
        detail = str(e)
        if "requires more system memory" in detail or "more system memory" in detail:
            response = (
                f"El modelo `{model_name}` es demasiado grande para la memoria disponible ahora. "
                f"Usa un modelo más ligero: `ollama pull {RECOMMENDED_OLLAMA_MODEL}`. "
                "Luego presiona `Actualizar modelos` y selecciona ese modelo en la barra lateral."
            )
        else:
            response = (
                f"No pude usar el modelo `{model_name}` en Ollama. "
                f"Verifica que esté instalado con `ollama list` o descárgalo con "
                f"`ollama pull {model_name}`. Detalle: {e}"
            )
    # Guardar en memoria si existe
    response = clean_answer_text(response)
    if memory is not None:
        add_memory_turn(memory, question, response)
    return response

def answer_from_brain(
    question,
    user_groups=None,
    memory=None,
    interaction_mode="unified",
    model_name=None,
    k=None,
    top_k=None,
    selected_sources=None,
    agency_enabled=True,
    show_sources=False,
    brain_profile=None,
    workspace_path="",
    quick_code_context="",
    notebooklm_enabled=False,
    notebooklm_active_id="",
    smart_search_enabled=False,
    tools_enabled=True,
    fast_mode=True,
    deep_thinking=False,
):
    direct_result = try_direct_technical_answer(question, memory)
    if direct_result:
        return direct_result

    if not UnifiedBrain or not BrainSourceContext:
        docs = retrieve(question, user_groups=user_groups, k=k, top_k=top_k, selected_sources=selected_sources)
        workspace_docs = retrieve_connected_workspace_docs(question, workspace_path) if workspace_path else []
        quick_code_docs = build_quick_code_docs(quick_code_context)
        notebooklm_docs = []
        if notebooklm_enabled:
            notebooklm_service = NotebookLMService(enabled=True, active_notebook_id=notebooklm_active_id)
            notebooklm_result = notebooklm_service.ask(question, notebooklm_active_id)
            if notebooklm_result.ok and notebooklm_result.answer:
                notebook_doc = notebooklm_result_to_doc(
                    notebooklm_result.to_dict(),
                    notebook_id=notebooklm_active_id,
                    question=question,
                )
                if notebook_doc:
                    notebooklm_docs.append(notebook_doc)
        docs = quick_code_docs + notebooklm_docs + docs + workspace_docs
        agency_matches = retrieve_agency_agents(question, limit=AGENCY_MATCH_LIMIT) if agency_enabled else []
        agency_context = build_agency_context(agency_matches, max_chars=AGENCY_CONTEXT_CHARS)
        brain_bundle = build_connected_brain_context(
            question,
            interaction_mode=interaction_mode,
            brain_profile=brain_profile or "unified",
            workspace_path=workspace_path,
            quick_code_context=quick_code_context,
        )
        answer = generate_answer(
            question,
            docs,
            memory,
            interaction_mode=interaction_mode,
            model_name=model_name,
            agency_context=agency_context,
            brain_context=brain_bundle["context"],
            show_sources=show_sources,
            fast_response=fast_mode,
        )
        return {
            "answer": answer,
            "sources": docs,
            "agency_agents": agency_matches,
            "brain_parts": brain_bundle["parts"],
            "notebooklm_used_count": len(notebooklm_docs),
            "mode": get_interaction_mode(interaction_mode)["label"],
            "source_contexts": [],
        }

    local_docs = []
    workspace_docs = []
    quick_code_docs = []
    notebooklm_docs = []
    if notebooklm_enabled:
        notebooklm_docs = []
    notebooklm_result = None
    notebooklm_message = ""
    agency_matches = []
    agency_context = ""
    brain_bundle = {"context": "", "profile": brain_profile or "unified", "parts": []}
    smart_search = None

    def local_sources_provider(message, options, route):
        nonlocal local_docs, workspace_docs, quick_code_docs
        contexts = []
        quick_code_docs = build_quick_code_docs(quick_code_context)
        if quick_code_docs:
            contexts.append(
                BrainSourceContext(
                    source="uploaded_files",
                    success=True,
                    confidence=0.95,
                    content="\n\n".join(
                        trim_prompt_text(doc.get("text", ""), MAX_DOC_CONTEXT_CHARS)
                        for doc in quick_code_docs
                    ),
                    metadata={"docs": quick_code_docs},
                )
            )

        local_docs = retrieve(message, user_groups=user_groups, k=k, top_k=top_k, selected_sources=selected_sources)
        workspace_docs = retrieve_connected_workspace_docs(message, workspace_path) if workspace_path else []
        docs = local_docs + workspace_docs
        if docs:
            contexts.append(
                BrainSourceContext(
                    source="local_sources",
                    success=True,
                    confidence=0.82,
                    content="\n\n".join(
                        f"[{doc.get('metadata', {}).get('title', doc.get('metadata', {}).get('source', 'fuente'))}] "
                        f"{trim_prompt_text(doc.get('text', ''), MAX_DOC_CONTEXT_CHARS)}"
                        for doc in docs[:6]
                    ),
                    metadata={"docs": docs},
                )
            )
        return contexts

    def notebooklm_provider(message, options, route):
        nonlocal notebooklm_result, notebooklm_docs, notebooklm_message
        service = NotebookLMService(enabled=True, active_notebook_id=notebooklm_active_id)
        notebooklm_result = service.ask(message, notebooklm_active_id)
        if notebooklm_result.ok and notebooklm_result.answer:
            notebook_doc = notebooklm_result_to_doc(
                notebooklm_result.to_dict(),
                notebook_id=notebooklm_active_id,
                question=message,
            )
            if notebook_doc:
                notebooklm_docs.append(notebook_doc)
                return BrainSourceContext(
                    source="notebooklm",
                    success=True,
                    confidence=0.8,
                    content=notebook_doc["text"],
                    references=notebooklm_result.references,
                    metadata={"docs": notebooklm_docs, "result": notebooklm_result.to_dict()},
                )
        notebooklm_message = notebooklm_status_message(notebooklm_result)
        return BrainSourceContext(source="notebooklm", success=False, error=notebooklm_message)

    def agency_provider(message, options, route):
        nonlocal agency_matches, agency_context
        agency_matches = retrieve_agency_agents(message, limit=AGENCY_MATCH_LIMIT) if agency_enabled else []
        agency_context = build_agency_context(agency_matches, max_chars=AGENCY_CONTEXT_CHARS)
        return BrainSourceContext(
            source="agency",
            success=bool(agency_context),
            confidence=0.7,
            content=agency_context,
            metadata={"agency_matches": agency_matches},
        )

    def openjarvis_provider(message, options, route):
        nonlocal brain_bundle
        brain_bundle = build_connected_brain_context(
            message,
            interaction_mode=interaction_mode,
            brain_profile=brain_profile or "unified",
            workspace_path=workspace_path,
            quick_code_context=quick_code_context,
        )
        return BrainSourceContext(
            source="openjarvis",
            success=bool(brain_bundle.get("context")),
            confidence=0.65,
            content=brain_bundle.get("context", ""),
            metadata={"brain_parts": brain_bundle.get("parts", []), "brain_profile": brain_bundle.get("profile")},
        )

    def web_provider(message, options, route):
        nonlocal smart_search
        smart_search = smart_web_search(message)
        content = smart_search.get("message", "")
        docs = []
        if smart_search.get("enabled") and content:
            docs.append(make_context_doc("web_search", "Busqueda inteligente", content, "web"))
        return BrainSourceContext(
            source="web_search",
            success=bool(smart_search.get("enabled")),
            confidence=0.55,
            content=content,
            error="" if smart_search.get("enabled") else content,
            metadata={"docs": docs, "smart_search": smart_search},
        )

    def tools_provider(message, options, route):
        content = build_tools_context(message, route)
        docs = [make_context_doc("tools", "Herramientas locales", content, "tool")] if content else []
        return BrainSourceContext(
            source="tools",
            success=bool(content),
            confidence=0.6,
            content=content,
            metadata={"docs": docs, "route": route},
        )

    unified = UnifiedBrain(
        providers={
            "local_sources": local_sources_provider,
            "notebooklm": notebooklm_provider,
            "agency": agency_provider,
            "openjarvis": openjarvis_provider,
            "web_search": web_provider,
            "tools": tools_provider,
        }
    )
    route = unified.route_question(question)
    options = {
        "fast_mode": fast_mode,
        "deep_thinking": deep_thinking,
        "web_search": smart_search_enabled,
        "notebooklm": notebooklm_enabled,
        "agency": agency_enabled,
        "openjarvis": True,
        "local_sources": True,
        "tools": tools_enabled,
    }
    contexts = unified.collect_context(question, options, route)
    merged_contexts = unified.merge_contexts(contexts)
    docs = brain_contexts_to_docs(merged_contexts)

    agency_contexts = brain_context_by_source(merged_contexts, "agency")
    if agency_contexts and not agency_context:
        agency_context = brain_context_content(agency_contexts[0])
        agency_matches = brain_context_metadata(agency_contexts[0]).get("agency_matches", [])

    openjarvis_contexts = brain_context_by_source(merged_contexts, "openjarvis")
    if openjarvis_contexts and not brain_bundle.get("context"):
        metadata = brain_context_metadata(openjarvis_contexts[0])
        brain_bundle = {
            "context": brain_context_content(openjarvis_contexts[0]),
            "profile": metadata.get("brain_profile") or brain_profile or "unified",
            "parts": metadata.get("brain_parts", []),
        }

    answer = generate_answer(
        question,
        docs,
        memory,
        interaction_mode=interaction_mode,
        model_name=model_name,
        agency_context=agency_context,
        brain_context=brain_bundle.get("context", ""),
        show_sources=show_sources,
        fast_response=fast_mode,
    )
    if smart_search_enabled and smart_search and not smart_search.get("enabled"):
        answer = f"{answer}\n\n{smart_search.get('message') or SMART_SEARCH_UNCONFIGURED_MESSAGE}"

    return {
        "answer": answer,
        "sources": docs,
        "agency_agents": agency_matches,
        "brain_parts": brain_bundle.get("parts", []),
        "notebooklm_used_count": len(notebooklm_docs),
        "notebooklm_message": notebooklm_message,
        "route": route,
        "source_contexts": [context.to_dict() for context in contexts],
        "mode": get_interaction_mode(interaction_mode)["label"],
    }

APP_NAME = "Libro de Programación con IA"
APP_SUBTITLE_DATE = "16 abr 2026"
SOURCE_LIST_LIMIT = int(os.getenv("TUTOR_IA_SOURCE_LIST_LIMIT", "20000"))
BANNER_IMAGE_URL = (
    "https://images.unsplash.com/photo-1515879218367-8466d910aaa4"
    "?auto=format&fit=crop&w=1600&q=80"
)

DEMO_SOURCES = [
    {"source": "demo:youtube-ia-app", "title": "Aprende a integrar IA a tu Aplicación", "type": "video", "fragments": 1},
    {"source": "demo:programar-20-min", "title": "Aprende a programar en 20 minutos", "type": "video", "fragments": 1},
    {"source": "demo:javascript-react", "title": "Aprendiendo Javascript y React desde cero", "type": "video", "fragments": 1},
    {"source": "demo:nextjs", "title": "Aprendiendo NextJS", "type": "video", "fragments": 1},
    {"source": "demo:clase-1-pdf", "title": "Apuntes-Curso-Desarrollo-IA-Clase-1.pdf", "type": "pdf", "fragments": 1},
    {"source": "demo:prompting", "title": "Artificial Intelligence Professional Prompting", "type": "docx", "fragments": 1},
    {"source": "demo:programar-2026", "title": "Así se programa en 2026: IA, agentes y automatización", "type": "video", "fragments": 1},
]

SOURCE_TYPE_LABELS = {
    "video": ("▶", "YouTube"),
    "pdf": ("PDF", "PDF"),
    "article": ("🌐", "Web"),
    "docx": ("DOC", "Documento"),
    "pptx": ("PPT", "Presentación"),
    "excel": ("XLS", "Hoja"),
    "text": ("TXT", "Nota"),
}

def inject_notebooklm_css():
    st.markdown(
        """
<style>
:root {
    --app-bg: #f5f7fb;
    --panel: #ffffff;
    --text: #1f2937;
    --muted: #6b7280;
    --blue: #2563eb;
    --border: #e5e7eb;
}
html, body, [data-testid="stAppViewContainer"] {
    background: var(--app-bg);
    color: var(--text);
    font-family: Inter, Poppins, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
}
[data-testid="stHeader"] { background: transparent; }
.block-container {
    max-width: 1240px;
    padding-top: 1rem;
    padding-bottom: 6rem;
}
section[data-testid="stSidebar"] {
    background: #eef3ff;
    border-right: 1px solid rgba(148, 163, 184, 0.24);
}
section[data-testid="stSidebar"] > div {
    padding: 0.85rem 0.9rem 2rem;
}
.nb-card {
    background: var(--panel);
    border: 1px solid var(--border);
    border-radius: 22px;
    box-shadow: 0 16px 40px rgba(15, 23, 42, 0.06);
}
.nb-sidebar-title, .nb-chat-topbar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    background: var(--panel);
    border: 1px solid var(--border);
    border-radius: 20px;
    padding: 16px 18px;
    margin-bottom: 12px;
    box-shadow: 0 12px 30px rgba(15, 23, 42, 0.05);
}
.nb-sidebar-title h2, .nb-chat-topbar h1 {
    margin: 0;
    font-size: 1.18rem;
    line-height: 1.2;
    font-weight: 700;
    color: #111827;
}
.nb-icon-actions {
    color: #334155;
    display: flex;
    gap: 12px;
    font-size: 1.15rem;
}
.nb-source-search {
    background: var(--panel);
    border: 1px solid var(--border);
    border-radius: 20px;
    padding: 13px 14px;
    color: var(--muted);
    margin: 10px 0 12px;
    box-shadow: 0 10px 22px rgba(15, 23, 42, 0.04);
}
.nb-mini-row {
    display: flex;
    gap: 8px;
    align-items: center;
    margin-top: 12px;
}
.nb-chip {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    border: 1px solid #dbe3ef;
    border-radius: 999px;
    padding: 6px 10px;
    color: #334155;
    background: #ffffff;
    font-size: 0.82rem;
}
.nb-source-count {
    color: var(--muted);
    font-size: 0.84rem;
    margin: 8px 0 4px;
}
.nb-banner {
    position: relative;
    min-height: 268px;
    border-radius: 22px;
    overflow: hidden;
    border: 1px solid rgba(15, 23, 42, 0.08);
    background-image:
        linear-gradient(90deg, rgba(2, 6, 23, 0.78), rgba(2, 6, 23, 0.30)),
        url('__BANNER_IMAGE__');
    background-size: cover;
    background-position: center;
    box-shadow: 0 22px 55px rgba(15, 23, 42, 0.16);
    margin: 10px 0 24px;
}
.nb-banner-content {
    position: absolute;
    left: 32px;
    right: 28px;
    bottom: 28px;
    color: #ffffff;
}
.nb-banner h2 {
    margin: 0 0 8px;
    font-size: clamp(2rem, 4vw, 3.1rem);
    line-height: 1.05;
    font-weight: 800;
    letter-spacing: 0;
}
.nb-banner p {
    margin: 0;
    color: rgba(255, 255, 255, 0.88);
    font-weight: 650;
}
.nb-personalize {
    position: absolute;
    right: 20px;
    top: 18px;
    border: 1px solid rgba(255,255,255,0.42);
    background: rgba(15, 23, 42, 0.58);
    color: #ffffff;
    border-radius: 999px;
    padding: 9px 14px;
    font-weight: 700;
    backdrop-filter: blur(10px);
}
.nb-description {
    background: #ffffff;
    border: 1px solid var(--border);
    border-radius: 22px;
    padding: 22px 24px;
    box-shadow: 0 14px 34px rgba(15, 23, 42, 0.05);
    margin-bottom: 14px;
}
.nb-description p {
    margin: 0 0 14px;
    font-size: 1.04rem;
    line-height: 1.65;
}
.nb-quote {
    color: #111827;
    font-weight: 700;
}
.nb-notes-label {
    display: inline-flex;
    align-items: center;
    gap: 8px;
    margin-top: 4px;
    color: #475569;
    font-weight: 650;
}
.nb-brain-panel {
    background: #ffffff;
    border: 1px solid var(--border);
    border-radius: 18px;
    padding: 18px 20px;
    box-shadow: 0 14px 34px rgba(15, 23, 42, 0.05);
    margin: 12px 0 16px;
}
.nb-brain-head {
    display: flex;
    align-items: flex-start;
    justify-content: space-between;
    gap: 12px;
    margin-bottom: 14px;
}
.nb-brain-head h3 {
    margin: 0;
    color: #111827;
    font-size: 1rem;
    line-height: 1.2;
}
.nb-brain-head p {
    margin: 5px 0 0;
    color: #64748b;
    font-size: 0.88rem;
}
.nb-brain-grid {
    display: grid;
    grid-template-columns: repeat(4, minmax(0, 1fr));
    gap: 10px;
}
.nb-brain-stat {
    border: 1px solid #e5e7eb;
    border-radius: 14px;
    padding: 12px;
    background: #f8fafc;
    min-height: 84px;
}
.nb-brain-stat strong {
    display: block;
    color: #0f172a;
    font-size: 1.1rem;
    line-height: 1.2;
    margin-bottom: 4px;
}
.nb-brain-stat span {
    display: block;
    color: #64748b;
    font-size: 0.8rem;
}
.nb-brain-pill {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    border-radius: 999px;
    padding: 6px 10px;
    background: #ecfeff;
    border: 1px solid #cffafe;
    color: #155e75;
    font-size: 0.8rem;
    font-weight: 700;
    white-space: nowrap;
}
.nb-tool-row {
    display: flex;
    flex-wrap: wrap;
    gap: 7px;
    margin-top: 12px;
}
.nb-tool-chip {
    border-radius: 999px;
    border: 1px solid #e2e8f0;
    background: #ffffff;
    color: #475569;
    padding: 5px 9px;
    font-size: 0.78rem;
}
.nb-input-meta {
    display: flex;
    justify-content: flex-end;
    color: var(--muted);
    font-size: 0.92rem;
    margin: 8px 4px 0;
}
div[data-testid="stChatMessage"] {
    background: #ffffff;
    border: 1px solid #e7edf5;
    border-radius: 18px;
    padding: 0.35rem 0.65rem;
    box-shadow: 0 10px 24px rgba(15, 23, 42, 0.04);
}
.stButton > button, .stDownloadButton > button {
    border-radius: 999px;
    border: 1px solid #d9e2ef;
    background: #ffffff;
    color: #1f2937;
    box-shadow: 0 8px 20px rgba(15, 23, 42, 0.04);
    transition: all 160ms ease;
}
.stButton > button:hover, .stDownloadButton > button:hover {
    border-color: var(--blue);
    color: var(--blue);
    transform: translateY(-1px);
}
[data-testid="stSidebar"] .stCheckbox {
    background: #ffffff;
    border: 1px solid transparent;
    border-radius: 14px;
    padding: 5px 8px;
    margin-bottom: 4px;
}
[data-testid="stSidebar"] .stCheckbox:hover {
    border-color: #dbeafe;
    background: #f8fbff;
}
@media (max-width: 900px) {
    .block-container { padding-left: 0.75rem; padding-right: 0.75rem; }
    .nb-banner { min-height: 230px; }
    .nb-banner-content { left: 20px; bottom: 22px; }
    .nb-brain-grid { grid-template-columns: repeat(2, minmax(0, 1fr)); }
}
</style>
""".replace("__BANNER_IMAGE__", BANNER_IMAGE_URL),
        unsafe_allow_html=True,
    )

def truncate_text(text, max_chars=48):
    text = str(text or "").strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"

def get_source_icon(source_type):
    return SOURCE_TYPE_LABELS.get(source_type, ("📄", "Documento"))[0]

def get_source_title(metadata):
    title = metadata.get("title") or metadata.get("source") or "Fuente sin título"
    return Path(title).name if isinstance(title, str) and "\\" in title else str(title)

@st.cache_data(ttl=30, show_spinner=False)
def get_source_summaries(user_groups=None):
    user_groups = normalize_groups(user_groups or ["public"])
    try:
        collection = get_collection()
        data = collection.get(include=["metadatas"], limit=SOURCE_LIST_LIMIT)
    except Exception:
        data = {"metadatas": []}

    sources = {}
    for metadata in data.get("metadatas") or []:
        if not metadata:
            continue
        group = metadata.get("access_group", "public")
        if "admin" not in user_groups and group not in user_groups:
            continue
        source = metadata.get("source") or metadata.get("title")
        if not source:
            continue
        if source not in sources:
            sources[source] = {
                "source": source,
                "title": get_source_title(metadata),
                "type": metadata.get("type", "docx"),
                "fragments": 0,
            }
        sources[source]["fragments"] += 1

    if not sources:
        return DEMO_SOURCES.copy()
    return sorted(sources.values(), key=lambda item: item["title"].lower())

def source_checkbox_key(source):
    digest = hashlib.sha1(str(source).encode("utf-8")).hexdigest()[:12]
    return f"source_active_{digest}"

def sync_all_source_checkboxes(source_keys):
    value = st.session_state.get("sources_select_all", True)
    for key in source_keys:
        st.session_state[key] = value

def render_sources_panel(source_summaries):
    st.markdown(
        """
<div class="nb-sidebar-title">
  <h2>Fuentes</h2>
  <div class="nb-icon-actions"><span>⊞</span></div>
</div>
""",
        unsafe_allow_html=True,
    )
    st.markdown(
        """
<div class="nb-source-search">
  🔎&nbsp;&nbsp;Buscar nuevas fuentes en la Web
  <div class="nb-mini-row">
    <span class="nb-chip">🌐 Web</span>
    <span class="nb-chip">✦ Fast Research</span>
    <span class="nb-chip">→</span>
  </div>
</div>
""",
        unsafe_allow_html=True,
    )

    source_keys = [source_checkbox_key(item["source"]) for item in source_summaries]
    if "sources_select_all" not in st.session_state:
        st.session_state.sources_select_all = True
    for key in source_keys:
        if key not in st.session_state:
            st.session_state[key] = st.session_state.sources_select_all

    st.checkbox(
        "Seleccionar todas las fuentes",
        key="sources_select_all",
        on_change=sync_all_source_checkboxes,
        args=(source_keys,),
    )

    selected_sources = []
    with st.container(height=420, border=False):
        for item in source_summaries:
            key = source_checkbox_key(item["source"])
            icon = get_source_icon(item.get("type"))
            title = truncate_text(item.get("title"), 42)
            if st.checkbox(f"{icon}  {title}", key=key):
                selected_sources.append(item["source"])

    st.markdown(
        f'<div class="nb-source-count">{len(selected_sources)}/{len(source_summaries)} fuentes activas</div>',
        unsafe_allow_html=True,
    )
    return selected_sources

def render_topbar():
    st.markdown(
        """
<div class="nb-chat-topbar">
  <h1>Chat</h1>
  <div class="nb-icon-actions"><span>☷</span><span>⚙</span><span>⋮</span></div>
</div>
""",
        unsafe_allow_html=True,
    )

def render_banner(source_count):
    subtitle = f"{source_count} fuentes · {APP_SUBTITLE_DATE}"
    st.markdown(
        f"""
<section class="nb-banner">
  <div class="nb-personalize">▣&nbsp; Personalizar</div>
  <div class="nb-banner-content">
    <h2>{escape(APP_NAME)}</h2>
    <p>{escape(subtitle)}</p>
  </div>
</section>
""",
        unsafe_allow_html=True,
    )

def render_intro_card():
    st.markdown(
        """
<div class="nb-description">
  <p>El libro es una guía técnica y estratégica diseñada para que los programadores evolucionen hacia el rol de Arquitectos de Orquestación en la era de la inteligencia artificial.</p>
  <p class="nb-quote">"A programar se aprende programando"</p>
  <div class="nb-notes-label">☷ Notas del creador</div>
</div>
""",
        unsafe_allow_html=True,
    )

def render_jarvis_brain_panel(source_count, active_source_count, agency_status, model_name, profile_key):
    summary = get_jarvis_stack_summary()
    profile = get_profile(profile_key)
    installed_models = get_installed_ollama_models()
    active_model = choose_llm_model(model_name, brain_context=build_unified_brain_context())
    model_display = f"Auto -> {active_model}" if model_name == AUTO_MODEL_OPTION else (active_model or model_name or "sin modelo")
    openjarvis = summary["openjarvis"]
    jarvis_mlx = summary["jarvis_mlx"]
    openjarvis_status = "Detectado" if openjarvis["available"] else "Sin ruta"
    voice_status = "STT/TTS listo" if jarvis_mlx["available"] and jarvis_mlx["speech_to_text"] else "Voz opcional"
    agency_count = agency_status.get("count", 0) if agency_status else 0
    profile_status = f"{summary['detected_profiles']} perfiles"
    base_tools = list(summary.get("tools", []))
    for integrated_tool in ["notebooklm", "agency", "openjarvis"]:
        if integrated_tool not in base_tools:
            base_tools.append(integrated_tool)
    tools = base_tools[:10]
    tool_html = "".join(f'<span class="nb-tool-chip">{escape(tool)}</span>' for tool in tools)
    if len(base_tools) > len(tools):
        tool_html += f'<span class="nb-tool-chip">+{len(base_tools) - len(tools)}</span>'

    st.markdown(
        f"""
<section class="nb-brain-panel">
  <div class="nb-brain-head">
    <div>
      <h3>Cerebro Unificado</h3>
      <p>{escape(profile["label"])} conectado a fuentes, Obsidian, Agency, OpenJarvis, NotebookLM y voz local.</p>
    </div>
    <span class="nb-brain-pill">local-first</span>
  </div>
  <div class="nb-brain-grid">
    <div class="nb-brain-stat"><strong>{escape(model_display)}</strong><span>{len(installed_models)} modelos Ollama</span></div>
    <div class="nb-brain-stat"><strong>{active_source_count}/{source_count}</strong><span>Fuentes activas</span></div>
    <div class="nb-brain-stat"><strong>{agency_count}</strong><span>Especialistas Agency</span></div>
    <div class="nb-brain-stat"><strong>{escape(profile_status)}</strong><span>{escape(openjarvis_status)} / {escape(voice_status)}</span></div>
  </div>
  <div class="nb-tool-row">{tool_html}</div>
</section>
""",
        unsafe_allow_html=True,
    )


def render_brain_connector_status_panel():
    connector = get_brain_connector()
    status = connector.get_status()
    local_status = status.get("local_brain", {})
    bridge_status = status.get("bridge_api", {})
    anthropic_status = status.get("anthropic", {})

    with st.expander("BrainConnector", expanded=False):
        st.caption(f"Cerebro local: {'conectado' if local_status.get('connected') else 'desconectado'}")
        st.caption(f"Bridge API: {'conectada' if bridge_status.get('connected') else 'no disponible'}")
        st.caption(f"Anthropic: {'configurado' if anthropic_status.get('configured') else 'no configurado'}")
        st.caption("Modo: local-first")
        st.code(local_status.get("root") or str(Path.cwd()), language="text")
        st.session_state.deep_thinking = st.checkbox(
            "Pensamiento profundo",
            value=st.session_state.get("deep_thinking", False),
            help="Usa mas contexto y Claude si esta configurado. Puede tardar mas.",
        )
        if st.button("Probar conexion del cerebro", key="test_brain_connector", use_container_width=True):
            health = connector.health_check()
            probe = connector.answer(
                "Responde en una frase: estas conectado?",
                options={
                    "local_first": True,
                    "fast_mode": True,
                    "bridge_api": True,
                    "anthropic": False,
                    "deep_thinking": False,
                },
            )
            local_ok = health.get("local_brain", {}).get("connected")
            bridge_ok = health.get("bridge_api", {}).get("connected")
            anthropic_ok = health.get("anthropic", {}).get("configured")
            st.write(f"Cerebro local: {'OK' if local_ok else 'no disponible'}")
            st.write(f"Bridge API: {'OK' if bridge_ok else 'no disponible'}")
            st.write(f"Anthropic: {'OK' if anthropic_ok else 'no configurado'}")
            st.write(f"Respuesta de prueba: {'OK' if probe.get('answer') else 'sin respuesta'}")


def transcribe_audio_prompt(audio_file):
    audio_bytes = audio_file.getvalue()
    if not audio_bytes:
        return ""

    digest = hashlib.sha1(audio_bytes).hexdigest()
    if (
        st.session_state.get("jarvis_voice_digest") == digest
        and st.session_state.get("jarvis_voice_text")
    ):
        return st.session_state.jarvis_voice_text

    suffix = Path(getattr(audio_file, "name", "") or "voice.wav").suffix or ".wav"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(audio_bytes)
            tmp_path = tmp.name

        segments, _ = get_whisper().transcribe(tmp_path, language="es")
        text = " ".join(segment.text.strip() for segment in segments if segment.text.strip())
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    st.session_state.jarvis_voice_digest = digest
    st.session_state.jarvis_voice_text = text
    return text


def render_voice_console():
    if not st.session_state.get("jarvis_voice_enabled", False):
        return None

    with st.expander("Entrada por voz Jarvis", expanded=False):
        if not hasattr(st, "audio_input"):
            st.info("Tu version de Streamlit no incluye entrada de microfono integrada.")
            return None

        audio_file = st.audio_input("Dictar pregunta", key="jarvis_voice_audio")
        if not audio_file:
            return None

        with st.spinner("Transcribiendo voz con Whisper local..."):
            text = transcribe_audio_prompt(audio_file)

        if not text:
            st.warning("No pude extraer texto claro del audio.")
            return None

        st.caption(text)
        if st.button("Enviar voz al cerebro", use_container_width=True):
            return text.strip()
    return None


# ============= STREAMLIT UI =============
st.set_page_config(
    page_title=APP_NAME,
    page_icon="📘",
    layout="wide",
    initial_sidebar_state="expanded",
)
inject_notebooklm_css()

# --- Control de sesión y autenticación ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.username = None
    st.session_state.user_groups = ["public"]
    st.session_state.memory = new_memory()
    st.session_state.interaction_mode = "unified"
    st.session_state.llm_model = AUTO_MODEL_OPTION
    st.session_state.agency_enabled = True
    st.session_state.jarvis_profile = "unified"
    st.session_state.jarvis_voice_enabled = False
    st.session_state.code_workspace_path = ""
    st.session_state.quick_code_context = ""
    st.session_state.notebooklm_enabled = NOTEBOOKLM_ENABLED_DEFAULT
    st.session_state.notebooklm_active_id = NOTEBOOKLM_ACTIVE_ID
    st.session_state.deep_thinking = False
if "agency_enabled" not in st.session_state:
    st.session_state.agency_enabled = True
if "jarvis_profile" not in st.session_state:
    st.session_state.jarvis_profile = "unified"
if "jarvis_voice_enabled" not in st.session_state:
    st.session_state.jarvis_voice_enabled = False
if "code_workspace_path" not in st.session_state:
    st.session_state.code_workspace_path = ""
if "quick_code_context" not in st.session_state:
    st.session_state.quick_code_context = ""
if "notebooklm_enabled" not in st.session_state:
    st.session_state.notebooklm_enabled = NOTEBOOKLM_ENABLED_DEFAULT
if "notebooklm_active_id" not in st.session_state:
    st.session_state.notebooklm_active_id = NOTEBOOKLM_ACTIVE_ID
if "deep_thinking" not in st.session_state:
    st.session_state.deep_thinking = False

if not st.session_state.logged_in:
    first_admin_mode = count_users() == 0
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Iniciar sesión")
        username = st.text_input("Usuario")
        password = st.text_input("Contraseña", type="password")
        if st.button("Entrar"):
            groups = check_user(username, password)
            if groups:
                st.session_state.logged_in = True
                st.session_state.username = normalize_username(username)
                st.session_state.user_groups = groups
                st.session_state.memory = new_memory()
                st.session_state.interaction_mode = "unified"
                st.session_state.llm_model = AUTO_MODEL_OPTION
                st.session_state.agency_enabled = True
                st.session_state.jarvis_profile = "unified"
                st.session_state.jarvis_voice_enabled = False
                st.session_state.code_workspace_path = ""
                st.session_state.quick_code_context = ""
                st.session_state.notebooklm_enabled = NOTEBOOKLM_ENABLED_DEFAULT
                st.session_state.notebooklm_active_id = NOTEBOOKLM_ACTIVE_ID
                st.session_state.deep_thinking = False
                st.rerun()
            else:
                st.error("Usuario o contraseña incorrectos")

    with col2:
        if first_admin_mode:
            st.subheader("Crear primer administrador")
            new_user = st.text_input("Usuario administrador")
            new_pass = st.text_input("Contraseña nueva", type="password")
            if st.button("Crear administrador"):
                ok, message = create_user(new_user, new_pass, "admin,public")
                if ok:
                    st.success("Administrador creado. Ahora inicia sesión.")
                else:
                    st.error(message)
        else:
            st.subheader("Registro cerrado")
            st.info("Un administrador puede crear nuevas cuentas desde el panel lateral.")
    st.stop()

# --- Interfaz principal tipo NotebookLM ---
allowed_groups = normalize_groups(st.session_state.user_groups)
source_summaries = get_source_summaries(allowed_groups)
using_demo_sources = all(str(item["source"]).startswith("demo:") for item in source_summaries)
display_source_count = 221 if using_demo_sources else len(source_summaries)

with st.sidebar:
    selected_sources = render_sources_panel(source_summaries)

    with st.expander("＋ Añadir fuentes", expanded=False):
        # Mantiene el control de acceso actual: cada fuente se guarda con su grupo.
        if "admin" in allowed_groups:
            access_options = sorted(set(get_all_groups()) | {"public", "admin"})
            access_group = st.selectbox("Grupo de acceso", access_options)
        else:
            access_group = st.selectbox("Grupo de acceso", allowed_groups)

        tab1, tab2, tab3 = st.tabs(["Archivo", "URL", "Texto"])

        with tab1:
            st.caption(f"Carga hasta {MAX_BATCH_SOURCES} archivos por lote.")
            uploaded_files = st.file_uploader(
                "Subir archivos",
                type=["pdf", "docx", "xlsx", "pptx"],
                accept_multiple_files=True,
            )
            if uploaded_files:
                st.caption(f"{len(uploaded_files)} archivos seleccionados.")

            if uploaded_files and st.button("Procesar archivos", key="process_files_batch"):
                files_to_process = uploaded_files[:MAX_BATCH_SOURCES]
                if len(uploaded_files) > MAX_BATCH_SOURCES:
                    st.warning(f"Se procesarán los primeros {MAX_BATCH_SOURCES} archivos de este lote.")

                progress = st.progress(0)
                status_box = st.empty()
                total = len(files_to_process)
                started_at = time.perf_counter()

                def on_extract(idx, total_files, source_name):
                    status_box.info(
                        f"Extrayendo archivo {idx}/{total_files}: {source_name} "
                        f"(máximo {FILE_EXTRACT_TIMEOUT}s por archivo)"
                    )
                    progress.progress(int((idx - 1) * 100 / total_files))

                def on_index(idx, total_files, source_name, done, total_fragments):
                    status_box.info(f"Guardando {source_name}: {done}/{total_fragments} fragmentos")
                    file_progress = 0.5 + (0.5 * done / max(total_fragments, 1))
                    progress.progress(int(((idx - 1) + file_progress) * 100 / total_files))

                def on_done(idx, total_files, source_name):
                    progress.progress(int(idx * 100 / total_files))

                results = ingest_uploaded_files_batch(
                    files_to_process,
                    access_group,
                    extraction_callback=on_extract,
                    index_callback=on_index,
                    done_callback=on_done,
                )
                progress.progress(100)
                elapsed = time.perf_counter() - started_at
                status_box.empty()
                get_source_summaries.clear()
                st.caption(f"Lote terminado en {elapsed:.1f} segundos.")
                render_ingest_results(results)

        with tab2:
            st.caption(f"Pega hasta {MAX_BATCH_SOURCES} URLs, una por línea.")
            urls_text = st.text_area(
                "URLs de videos o artículos",
                height=180,
                placeholder="https://www.youtube.com/watch?v=...\nhttps://ejemplo.com/articulo",
            )
            transcribe_audio = st.checkbox(
                "Transcribir audio si el video no tiene subtítulos",
                value=False,
                help="Usa Whisper local. Puede tardar varios minutos por video.",
            )
            urls = parse_url_list(urls_text)
            if urls:
                st.caption(f"{len(urls)} URLs detectadas.")

            if st.button("Procesar URLs", key="process_urls_batch"):
                if not urls:
                    st.warning("Pega al menos una URL válida para procesar.")
                else:
                    urls_to_process = urls[:MAX_BATCH_SOURCES]
                    if len(urls) > MAX_BATCH_SOURCES:
                        st.warning(f"Se procesarán las primeras {MAX_BATCH_SOURCES} URLs de este lote.")

                    progress = st.progress(0)
                    status_box = st.empty()
                    total = len(urls_to_process)
                    def on_url_progress(done, total_urls, url, phase):
                        status_box.info(
                            f"{phase.title()} URL {min(done + 1, total_urls)}/{total_urls}: {url} "
                            f"(maximo {URL_EXTRACT_TIMEOUT}s por URL)"
                        )
                        progress.progress(int(done * 100 / max(total_urls, 1)))

                    results = ingest_url_sources_batch(
                        urls_to_process,
                        access_group,
                        audio_fallback=transcribe_audio,
                        progress_callback=on_url_progress,
                    )
                    progress = progress
                    status_box = status_box
                    total = total
                    for idx, url in enumerate([], start=1):
                        status_box.info(
                            f"Procesando URL {idx}/{total}: {url} "
                            f"(máximo {URL_EXTRACT_TIMEOUT}s por URL)"
                        )
                        results.append(ingest_url_source(url, access_group, audio_fallback=transcribe_audio))
                        progress.progress(int(idx * 100 / total))

                    progress.progress(100)
                    status_box.empty()
                    get_source_summaries.clear()
                    render_ingest_results(results)

        with tab3:
            manual_text = st.text_area("Texto libre")
            title = st.text_input("Título", "Nota rápida")
            tags_str = st.text_input("Etiquetas (opcional)")
            if st.button("Añadir texto"):
                if manual_text.strip():
                    tags = [t.strip() for t in tags_str.split(",") if t.strip()] if tags_str else None
                    raw = process_manual_text(manual_text.strip(), title.strip() or "Nota manual", tags)
                    stats = add_to_db(raw, access_group)
                    get_source_summaries.clear()
                    st.success(f"Texto añadido ({stats['indexed']} fragmentos nuevos)")
                    if stats["skipped"]:
                        st.info(f"{stats['skipped']} fragmentos ya existían y se saltaron.")
                else:
                    st.warning("Escribe algún texto para poder indexarlo.")

    with st.expander("Cuenta, modelo y acceso", expanded=False):
        st.caption(f"{st.session_state.username} · grupos: {', '.join(st.session_state.user_groups)}")
        col_a, col_b = st.columns(2)
        with col_a:
            if st.button("Limpiar chat", use_container_width=True):
                st.session_state.memory = new_memory()
                st.rerun()
        with col_b:
            if st.button("Salir", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.username = None
                st.session_state.user_groups = ["public"]
                st.session_state.memory = new_memory()
                st.session_state.interaction_mode = "unified"
                st.session_state.llm_model = AUTO_MODEL_OPTION
                st.session_state.agency_enabled = True
                st.session_state.jarvis_profile = "unified"
                st.session_state.jarvis_voice_enabled = False
                st.session_state.code_workspace_path = ""
                st.session_state.quick_code_context = ""
                st.session_state.notebooklm_enabled = NOTEBOOKLM_ENABLED_DEFAULT
                st.session_state.notebooklm_active_id = NOTEBOOKLM_ACTIVE_ID
                st.session_state.deep_thinking = False
                st.rerun()

        st.divider()
        st.subheader("Modelo local")
        if st.button("Actualizar modelos", key="refresh_ollama_models"):
            get_installed_ollama_models.clear()
            st.rerun()

        installed_models = get_installed_ollama_models()
        if installed_models:
            installed_models = sort_ollama_models(installed_models)
            model_options = [AUTO_MODEL_OPTION] + installed_models
            current_selection = st.session_state.get("llm_model", AUTO_MODEL_OPTION)
            model_index = model_options.index(current_selection) if current_selection in model_options else 0
            st.session_state.llm_model = st.selectbox("Modelo de Ollama", model_options, index=model_index)
            model_plan = get_model_plan(installed_models)
            st.caption(
                "Auto usa "
                f"{model_plan.get('fast') or 'sin rapido'} para tareas ligeras y "
                f"{model_plan.get('code') or 'sin modelo de codigo'} para programacion."
            )
            missing_models = model_plan.get("recommended_pulls", [])
            if missing_models:
                with st.expander("Modelos locales recomendados", expanded=False):
                    for item in missing_models[:3]:
                        st.caption(f"{item['model']} - {item['use']}")
                        st.code(item["command"], language="powershell")
        else:
            st.session_state.llm_model = None
            st.warning("No hay modelos de Ollama instalados.")
            st.code(f"ollama pull {RECOMMENDED_OLLAMA_MODEL}", language="powershell")

        st.divider()
        st.subheader("Cerebro conectado")
        render_brain_connector_status_panel()
        agency_status = get_agency_status()
        if agency_status["available"]:
            st.session_state.agency_enabled = True
            st.caption(
                f"Agency activo: {agency_status['count']} especialistas conectados."
            )
            if st.button("Recargar Agency", key="reload_agency"):
                clear_agency_cache()
                st.rerun()
        else:
            st.session_state.agency_enabled = False
            st.warning(agency_status["message"])

        st.session_state.jarvis_profile = "unified"
        jarvis_summary = get_jarvis_stack_summary()
        profile_count = len(jarvis_summary.get("profiles", []))
        tool_count = len(jarvis_summary.get("tools", []))
        st.caption(f"OpenJarvis conectado: {profile_count} perfiles internos, {tool_count} herramientas conceptuales.")
        st.session_state.jarvis_voice_enabled = st.checkbox(
            "Entrada por voz local",
            value=st.session_state.get("jarvis_voice_enabled", False),
            help="Usa el patron de Jarvis MLX con Whisper local cuando Streamlit permite microfono.",
        )

        with st.expander("Cerebro NotebookLM", expanded=False):
            st.session_state.notebooklm_enabled = st.checkbox(
                "Cerebro NotebookLM",
                value=st.session_state.get("notebooklm_enabled", NOTEBOOKLM_ENABLED_DEFAULT),
                help="Consulta Google NotebookLM como una capa adicional del cerebro. Requiere notebooklm login.",
            )
            notebooklm_service = get_notebooklm_service(enabled=st.session_state.notebooklm_enabled)
            notebooklm_status = notebooklm_service.get_status(check_auth=False)
            if st.session_state.notebooklm_enabled:
                st.caption("Cerebro NotebookLM activado.")
            else:
                st.caption(notebooklm_status["message"])

            if st.button("Actualizar notebooks", key="refresh_notebooklm_notebooks", use_container_width=True):
                result = notebooklm_service.list_notebooks()
                if result.ok:
                    st.session_state.notebooklm_notebooks = result.data or []
                    st.success(f"{len(st.session_state.notebooklm_notebooks)} notebooks encontrados.")
                else:
                    st.session_state.notebooklm_notebooks = []
                    st.warning(result.message)

            notebooks = st.session_state.get("notebooklm_notebooks", [])
            if notebooks:
                labels = [f"{item.get('title') or 'Notebook'} ({item.get('id')})" for item in notebooks]
                ids = [item.get("id", "") for item in notebooks]
                current_id = st.session_state.get("notebooklm_active_id", "")
                current_index = ids.index(current_id) if current_id in ids else 0
                selected_label = st.selectbox("Notebook activo", labels, index=current_index)
                st.session_state.notebooklm_active_id = ids[labels.index(selected_label)]
            else:
                st.session_state.notebooklm_active_id = st.text_input(
                    "ID de notebook activo",
                    value=st.session_state.get("notebooklm_active_id", ""),
                    placeholder="Pega aqui el ID del notebook",
                    key="notebooklm_active_id_input",
                ).strip()

            notebook_url = st.text_input("Agregar URL a NotebookLM", key="notebooklm_url_input")
            if st.button("Agregar URL", key="notebooklm_add_url", use_container_width=True):
                result = notebooklm_service.add_url_source(st.session_state.notebooklm_active_id, notebook_url)
                if result.ok:
                    st.success(result.message)
                else:
                    st.warning(result.message)

            notebook_file = st.file_uploader(
                "Agregar archivo al notebook",
                type=["pdf", "txt", "md", "docx", "csv", "png", "jpg", "jpeg", "webp", "mp3", "mp4", "wav", "m4a"],
                key="notebooklm_file_uploader",
            )
            if notebook_file and st.button("Agregar archivo", key="notebooklm_add_file", use_container_width=True):
                suffix = Path(notebook_file.name).suffix
                temp_path = ""
                try:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                        temp_file.write(bytes(notebook_file.getbuffer()))
                        temp_path = temp_file.name
                    result = notebooklm_service.add_file_source(st.session_state.notebooklm_active_id, temp_path)
                    if result.ok:
                        st.success(result.message)
                    else:
                        st.warning(result.message)
                finally:
                    if temp_path:
                        try:
                            Path(temp_path).unlink(missing_ok=True)
                        except OSError:
                            pass

            notebook_question = st.text_input(
                "Preguntar solo a NotebookLM",
                placeholder="Pregunta puntual sobre el notebook activo",
                key="notebooklm_direct_question",
            )
            if st.button("Preguntar al notebook", key="notebooklm_direct_ask", use_container_width=True):
                result = notebooklm_service.ask(notebook_question, st.session_state.notebooklm_active_id)
                if result.ok:
                    st.info(result.answer or "NotebookLM no devolvio texto.")
                else:
                    st.warning(result.message)

        with st.expander("Proyecto de programacion", expanded=False):
            st.caption("Conecta una carpeta de codigo para revisar, mejorar o disenar software completo con contexto real.")
            st.session_state.code_workspace_path = st.text_input(
                "Carpeta del proyecto",
                value=st.session_state.get("code_workspace_path", ""),
                placeholder="ruta/al/proyecto",
                key="code_workspace_path_input",
            ).strip()
            workspace_summary = summarize_workspace(st.session_state.code_workspace_path)
            if workspace_summary.get("available"):
                st.success(workspace_summary["message"])
                language_summary = ", ".join(
                    f"{key}:{value}"
                    for key, value in sorted(workspace_summary.get("languages", {}).items())
                )
                if language_summary:
                    st.caption(language_summary)
            elif st.session_state.code_workspace_path:
                st.warning(workspace_summary["message"])

            st.session_state.quick_code_context = st.text_area(
                "Codigo, error o requerimiento rapido",
                value=st.session_state.get("quick_code_context", ""),
                height=150,
                placeholder="Pega aqui un fragmento de codigo, traceback, esquema SQL o requisitos del proyecto.",
                key="quick_code_context_input",
            )

            skill_labels = ", ".join(skill["label"] for skill in get_programming_skill_catalog())
            st.caption(f"Habilidades internas: {skill_labels}")

        st.divider()
        if st.button("Estado de la base"):
            st.write(f"Fragmentos indexados: {get_collection().count()}")

        if "admin" in allowed_groups:
            st.divider()
            st.subheader("Gestionar usuarios")
            admin_new_user = st.text_input("Nuevo usuario", key="admin_new_user")
            admin_new_pass = st.text_input("Nueva contraseña", type="password", key="admin_new_pass")
            admin_groups = st.text_input("Grupos", "public", key="admin_groups")
            if st.button("Crear usuario", key="admin_create_user"):
                ok, message = create_user(admin_new_user, admin_new_pass, admin_groups)
                if ok:
                    st.success(message)
                else:
                    st.error(message)

active_source_count = len(selected_sources)

render_topbar()
render_banner(display_source_count)
render_intro_card()
render_jarvis_brain_panel(
    len(source_summaries),
    active_source_count,
    agency_status,
    st.session_state.get("llm_model"),
    st.session_state.get("jarvis_profile"),
)

action_cols = st.columns([1.25, 0.7, 0.7, 0.8, 3])
with action_cols[0]:
    if st.button("⚑ Guardar en una nota", use_container_width=True):
        st.toast("Nota guardada en esta sesión.")
with action_cols[1]:
    if st.button("Copiar", use_container_width=True):
        st.toast("Texto listo para copiar desde el navegador.")
with action_cols[2]:
    if st.button("👍", help="Me gusta", use_container_width=True):
        st.toast("Gracias por tu valoración.")
with action_cols[3]:
    if st.button("👎", help="No me gusta", use_container_width=True):
        st.toast("Lo tendré en cuenta para mejorar.")

selected_mode_key = "unified"

st.session_state.interaction_mode = selected_mode_key
selected_mode = get_interaction_mode(selected_mode_key)
st.caption(selected_mode["description"])

quick_question = None
quick_cols = st.columns(3)
for idx, prompt_text in enumerate(selected_mode["quick_prompts"]):
    with quick_cols[idx]:
        if st.button(prompt_text, key=f"quick_{selected_mode_key}_{idx}", use_container_width=True):
            quick_question = prompt_text

voice_question = render_voice_console()

st.divider()
for msg in st.session_state.memory:
    if msg.get("role") == "human":
        with st.chat_message("user"):
            st.markdown(msg.get("content", ""))
    elif msg.get("role") == "ai":
        with st.chat_message("assistant"):
            st.markdown(msg.get("content", ""))

use_agency_for_question = bool(st.session_state.get("agency_enabled", True) and agency_status.get("available"))
use_jarvis_for_question = True
use_notebooklm_for_question = bool(st.session_state.get("notebooklm_enabled", False))

if not selected_sources:
    if use_agency_for_question or use_jarvis_for_question:
        st.info("No hay fuentes activas; el cerebro unificado puede responder con sus capas internas cuando aplique.")
    else:
        st.warning("Selecciona al menos una fuente para que el cerebro pueda responder con contexto.")

input_meta = f"{active_source_count} fuentes"
if use_agency_for_question:
    input_meta += " + Agency"
input_meta += " + OpenJarvis"
if st.session_state.get("code_workspace_path"):
    input_meta += " + Proyecto"
if st.session_state.get("quick_code_context", "").strip():
    input_meta += " + Codigo"
if use_notebooklm_for_question:
    input_meta += " + NotebookLM"
if st.session_state.get("jarvis_voice_enabled", False):
    input_meta += " + Voz"
if st.session_state.get("deep_thinking", False):
    input_meta += " + Profundo"
st.markdown(
    f'<div class="nb-input-meta">{input_meta}</div>',
    unsafe_allow_html=True,
)
chat_question = st.chat_input("Mensaje a tu asistente de programacion...")
question = (quick_question or voice_question or chat_question or "").strip()
if question:
    with st.chat_message("user"):
        st.markdown(question)
    with st.chat_message("assistant"):
        with st.spinner("Consultando Cerebro Unificado..."):
            deep_thinking = bool(st.session_state.get("deep_thinking", False))
            result = try_direct_technical_answer(question, st.session_state.memory)
            if not result:
                connector_options = {
                    "local_first": True,
                    "fast_mode": not deep_thinking,
                    "bridge_api": True,
                    "anthropic": True,
                    "deep_thinking": deep_thinking,
                    "response_profile": "deep" if deep_thinking else "web_fast",
                    "session_id": f"streamlit:{st.session_state.username or 'default'}",
                    "interaction_mode": st.session_state.interaction_mode,
                    "mode": st.session_state.interaction_mode,
                    "model": st.session_state.get("llm_model"),
                    "selected_sources": selected_sources,
                    "agency_enabled": use_agency_for_question,
                    "show_sources": source_requested(question),
                    "jarvis_profile": st.session_state.get("jarvis_profile"),
                    "workspace_path": st.session_state.get("code_workspace_path", ""),
                    "quick_code_context": st.session_state.get("quick_code_context", "").strip(),
                    "notebookLM": use_notebooklm_for_question,
                    "notebooklm_enabled": use_notebooklm_for_question,
                    "notebooklm_active_id": st.session_state.get("notebooklm_active_id", ""),
                    "smart_search": bool(st.session_state.get("smart_search_enabled", False)),
                    "tools": True,
                    "history": st.session_state.memory,
                    "memory": st.session_state.memory,
                    "max_tokens": 3500,
                    "max_output_tokens": 3500,
                    "max_new_tokens": 3500,
                    "num_predict": 4096,
                    "temperature": 0.25,
                    "top_p": 0.9,
                }
                result = get_brain_connector().answer(question, options=connector_options)
                connector_sources = set(result.get("sources_used", []))
                trusted_connector_sources = {
                    "bridge_api",
                    "anthropic",
                    "database_generator",
                    "technical_generators",
                    "code_interpreter_router",
                    "conversation_resolver",
                    "fallback_programming_answer",
                }
                use_streamlit_local_fallback = not result.get("answer") or not connector_sources.intersection(trusted_connector_sources)
                if use_streamlit_local_fallback:
                    legacy_result = answer_from_brain(
                        question,
                        user_groups=st.session_state.user_groups,
                        memory=st.session_state.memory,
                        interaction_mode=st.session_state.interaction_mode,
                        model_name=st.session_state.get("llm_model"),
                        selected_sources=selected_sources,
                        agency_enabled=use_agency_for_question,
                        show_sources=source_requested(question),
                        brain_profile=st.session_state.get("jarvis_profile"),
                        workspace_path=st.session_state.get("code_workspace_path", ""),
                        quick_code_context=st.session_state.get("quick_code_context", "").strip(),
                        notebooklm_enabled=use_notebooklm_for_question,
                        notebooklm_active_id=st.session_state.get("notebooklm_active_id", ""),
                        smart_search_enabled=bool(st.session_state.get("smart_search_enabled", False)),
                        tools_enabled=True,
                        fast_mode=not deep_thinking,
                        deep_thinking=deep_thinking,
                    )
                    legacy_result["connector_result"] = result
                    legacy_result["sources_used"] = ["streamlit_local_fallback"]
                    legacy_result["success"] = True
                    result = legacy_result
                else:
                    add_memory_turn(st.session_state.memory, question, result.get("answer", ""))
            answer = result["answer"]
            docs = result.get("sources", [])
            agency_matches = result.get("agency_agents", [])
            st.markdown(answer)
            if result.get("notebooklm_message") and use_notebooklm_for_question and not result.get("notebooklm_used_count"):
                st.warning(result["notebooklm_message"])
            if docs:
                with st.expander("Fuentes usadas"):
                    for d in docs:
                        meta = d["metadata"]
                        st.caption(f"{meta.get('type','?')}: {meta.get('title', meta.get('source',''))}")
            else:
                source_labels = [
                    item.get("source", "")
                    for item in result.get("source_results", [])
                    if item.get("success")
                ]
                if source_labels:
                    st.caption(f"Fuentes usadas: {', '.join(source_labels)}")
                else:
                    st.caption("Sin fuentes recuperadas.")
            if agency_matches:
                with st.expander("Especialistas Agency usados"):
                    for agent in agency_matches:
                        st.caption(
                            f"{agent.get('name')} ({agent.get('category')}) - {agent.get('relative_path')}"
                        )

st.stop()

