import json
import sys
from pathlib import Path

import fitz
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation


def extract_pdf(file_path, source_name):
    chunks = []
    with fitz.open(file_path) as doc:
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                chunks.append({
                    "text": text,
                    "metadata": {"source": source_name, "type": "pdf", "page": i + 1},
                })
    return chunks


def extract_docx(file_path, source_name):
    doc = DocxDocument(file_path)
    full = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            full.append(paragraph.text)
    for table in doc.tables:
        rows = [" | ".join(cell.text for cell in row.cells) for row in table.rows]
        full.append("\n".join(rows))
    text = "\n\n".join(full)
    return [{"text": text, "metadata": {"source": source_name, "type": "docx"}}] if text.strip() else []


def extract_xlsx(file_path, source_name):
    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    chunks = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(" | ".join(str(cell) if cell is not None else "" for cell in row))
        if rows:
            chunks.append({
                "text": f"Hoja: {name}\n" + "\n".join(rows),
                "metadata": {"source": source_name, "type": "excel", "sheet": name},
            })
    wb.close()
    return chunks


def extract_pptx(file_path, source_name):
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(paragraph.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Nota]: {notes}")
        if texts:
            chunks.append({
                "text": f"Diapositiva {i + 1}\n" + "\n".join(texts),
                "metadata": {"source": source_name, "type": "pptx", "slide": i + 1},
            })
    return chunks


def main():
    if len(sys.argv) != 4:
        raise SystemExit("Uso: extract_worker.py <archivo> <nombre_fuente> <salida_json>")

    file_path = Path(sys.argv[1])
    source_name = sys.argv[2]
    out_path = Path(sys.argv[3])
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        chunks = extract_pdf(file_path, source_name)
    elif ext == ".docx":
        chunks = extract_docx(file_path, source_name)
    elif ext == ".xlsx":
        chunks = extract_xlsx(file_path, source_name)
    elif ext == ".pptx":
        chunks = extract_pptx(file_path, source_name)
    else:
        chunks = []

    out_path.write_text(json.dumps(chunks, ensure_ascii=False), encoding="utf-8")


if __name__ == "__main__":
    main()
